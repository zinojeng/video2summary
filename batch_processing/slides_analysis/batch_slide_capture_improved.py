#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
批量使用改進模式捕獲幻燈片的工具
可處理單支影片或整個資料夾中的影片 (支援遞迴)
"""

import os
import sys
import argparse
import time
import json
import shutil
from pathlib import Path
from typing import List, Dict, Optional, Tuple

from improved_slide_capture import capture_slides_improved


class BatchImprovedSlideCapture:
    """批量改進幻燈片捕獲器"""

    def __init__(
        self,
        threshold: float = 0.85,
        auto_select: bool = False,
        recursive: bool = False,
        force: bool = False,
        list_only: bool = False,
        yes: bool = False,
        generate_ppt: bool = True,
    ):
        self.threshold = threshold
        self.auto_select = auto_select
        self.recursive = recursive
        self.force = force
        self.list_only = list_only
        self.yes = yes
        self.generate_ppt = generate_ppt
        self.video_extensions = {".mp4",
                                  ".avi",
                                  ".mkv",
                                  ".mov",
                                  ".wmv",
                                  ".flv",
                                  ".webm"}
        self.processed_videos: List[Dict] = []
        self.failed_videos: List[Dict] = []
        self.stats = {
            "total_videos": 0,
            "processed": 0,
            "failed": 0,
            "total_slides": 0,
            "total_groups": 0,
            "processing_time": 0.0,
        }

    def find_videos(self, path: str) -> List[str]:
        """尋找需要處理的影片"""
        video_files: List[str] = []
        target = Path(path)

        if target.is_file():
            if target.suffix.lower() in self.video_extensions and not target.name.startswith("._"):
                video_files.append(str(target))
            else:
                print(f"警告：{path} 不是支援的影片格式或為系統檔案")
        elif target.is_dir():
            iterator = target.rglob if self.recursive else target.glob
            for ext in self.video_extensions:
                for candidate in iterator(f"*{ext}"):
                    if candidate.is_file() and not candidate.name.startswith("._"):
                        video_files.append(str(candidate))
        else:
            print(f"錯誤：{path} 不存在")

        return sorted(set(video_files))

    def process_video(self, video_path: str) -> Dict:
        """處理單支影片"""
        print(f"\n{'=' * 60}")
        print(f"處理影片: {os.path.basename(video_path)}")
        print(f"路徑: {video_path}")
        print(f"{'=' * 60}")

        video_dir = os.path.dirname(video_path)
        video_name = os.path.splitext(os.path.basename(video_path))[0]
        output_folder = os.path.join(video_dir, f"slides_improved_{video_name}")

        if not self.force:
            possible_folders = [
                f"slides_improved_{video_name}",
                f"{video_name}_slides_improved",
                f"{video_name}_slides",
                f"{video_name}_slide",
                "slides_improved",
                "slides",
                "slide",
                "Slides",
                "Slide",
            ]
            for folder_name in possible_folders:
                folder_path = os.path.join(video_dir, folder_name)
                if os.path.isdir(folder_path):
                    try:
                        has_images = any(
                            f.lower().endswith((".jpg", ".jpeg", ".png", ".gif", ".webp"))
                            for f in os.listdir(folder_path)
                            if not f.startswith("._")
                        )
                    except OSError:
                        has_images = False
                    if has_images:
                        print(f"⚠️  已存在幻燈片資料夾 '{folder_name}'，跳過處理")
                        return {
                            "status": "skipped",
                            "video": video_path,
                            "output": folder_path,
                            "reason": "already_processed",
                            "existing_folder": folder_name,
                        }

        start = time.time()
        try:
            print("使用改進模式捕獲幻燈片...")
            print(f"相似度閾值: {self.threshold}")
            success, result = capture_slides_improved(video_path, output_folder, self.threshold)

            if not success:
                print(f"\n❌ 捕獲失敗: {result.get('error', '未知錯誤')}")
                return {
                    "status": "failed",
                    "video": video_path,
                    "error": result.get("error", "未知錯誤"),
                }

            elapsed = time.time() - start
            metadata = self._load_metadata(output_folder)
            group_count = self._calculate_group_count(metadata)

            self.stats["total_slides"] += result.get("slide_count", 0)
            self.stats["total_groups"] += group_count

            ppt_path: Optional[str] = None
            ppt_error: Optional[str] = None
            if self.generate_ppt and result.get("slide_count", 0) > 0:
                ppt_success, ppt_output = self._generate_powerpoint(output_folder)
                if ppt_success:
                    ppt_path = ppt_output
                else:
                    ppt_error = ppt_output

            print("\n✅ 捕獲成功！")
            print(f"   幻燈片數: {result.get('slide_count', 0)}")
            print(f"   推估分組數: {group_count}")
            print(f"   耗時: {elapsed:.1f} 秒")
            print(f"   輸出資料夾: {output_folder}")
            if ppt_path:
                print(f"   PowerPoint: {ppt_path}")
            elif ppt_error:
                print(f"   PowerPoint 生成失敗: {ppt_error}")

            if self.auto_select and result.get("slide_count", 0) > 0:
                print("\n自動挑選代表幻燈片...")
                try:
                    selected_folder = self._auto_select_best(output_folder, metadata)
                    print(f"   已輸出到: {selected_folder}")
                except Exception as exc:
                    print(f"   自動挑選時發生錯誤: {exc}")

            try:
                self._create_simple_preview(output_folder, metadata)
                print("\n生成簡易預覽: preview.html")
            except Exception as exc:
                print(f"生成預覽時發生錯誤: {exc}")

            return {
                "status": "success",
                "video": video_path,
                "output": output_folder,
                "slide_count": result.get("slide_count", 0),
                "group_count": group_count,
                "processing_time": elapsed,
                "ppt_file": ppt_path,
            }
        except Exception as exc:
            print(f"\n❌ 處理時發生錯誤: {exc}")
            return {
                "status": "failed",
                "video": video_path,
                "error": str(exc),
            }

    def _load_metadata(self, slides_folder: str) -> Dict:
        metadata_path = os.path.join(slides_folder, "slides_metadata.json")
        if not os.path.exists(metadata_path):
            raise FileNotFoundError(f"找不到元數據檔案: {metadata_path}")
        with open(metadata_path, "r", encoding="utf-8") as handle:
            return json.load(handle)

    def _calculate_group_count(self, metadata: Dict) -> int:
        group_ids = set()
        singles = 0
        for slide in metadata.get("slides", []):
            group_id = slide.get("group_id", -1)
            if group_id is None or group_id < 0:
                singles += 1
            else:
                group_ids.add(group_id)
        return len(group_ids) + singles

    def _generate_powerpoint(self, slides_folder: str) -> Tuple[bool, str]:
        """將幻燈片影像生成 PowerPoint 檔案"""
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            return False, "需要安裝 python-pptx 套件 (pip install python-pptx)"

        if not os.path.isdir(slides_folder):
            return False, "找不到幻燈片資料夾"

        image_files = []
        for filename in sorted(os.listdir(slides_folder)):
            if filename.startswith("._"):
                continue
            if filename.lower().endswith((".jpg", ".jpeg", ".png")):
                image_files.append(os.path.join(slides_folder, filename))

        if not image_files:
            return False, "幻燈片資料夾內沒有圖片"

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        blank_layout = prs.slide_layouts[6]

        for image_path in image_files:
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=slide_width, height=slide_height)

        output_file = os.path.join(
            os.path.dirname(slides_folder),
            f"{os.path.basename(slides_folder)}.pptx"
        )

        try:
            prs.save(output_file)
        except Exception as err:
            return False, str(err)

        return True, output_file

    def _auto_select_best(self, slides_folder: str, metadata: Dict) -> str:
        selected_folder = os.path.join(slides_folder, "selected_slides")
        os.makedirs(selected_folder, exist_ok=True)

        grouped: Dict[int, List[Dict]] = {}
        singles: List[Dict] = []
        for slide in metadata.get("slides", []):
            group_id = slide.get("group_id", -1)
            if group_id is None or group_id < 0:
                singles.append(slide)
            else:
                grouped.setdefault(group_id, []).append(slide)

        copied = 0
        for slide in singles:
            source = os.path.join(slides_folder, slide["filename"])
            target = os.path.join(selected_folder, slide["filename"])
            if os.path.exists(source):
                shutil.copy2(source, target)
                copied += 1

        for group_id, slides in grouped.items():
            # 改進模式已為每組挑選最佳幀，因此直接選第一張
            chosen = slides[0]
            source = os.path.join(slides_folder, chosen["filename"])
            target_name = f"group_{group_id:03d}.jpg"
            target = os.path.join(selected_folder, target_name)
            if os.path.exists(source):
                shutil.copy2(source, target)
                copied += 1

        if copied == 0:
            print("   沒有可複製的幻燈片")
        else:
            print(f"   已複製 {copied} 張幻燈片")
        return selected_folder

    def _create_simple_preview(self, slides_folder: str, metadata: Dict):
        fps = metadata.get("fps", 0) or 0
        total_frames = metadata.get("total_frames", 0) or 0
        duration = total_frames / fps if fps else 0

        html_lines = [
            "<!DOCTYPE html>",
            "<html>",
            "<head>",
            "    <meta charset=\"UTF-8\">",
            f"    <title>改進幻燈片預覽 - {os.path.basename(slides_folder)}</title>",
            "    <style>",
            "        body { font-family: Arial, sans-serif; margin: 20px; background-color: #f5f5f5; }",
            "        .summary, .slide { background: #fff; border-radius: 8px; padding: 16px; margin-bottom: 16px;",
            "            box-shadow: 0 2px 4px rgba(0,0,0,0.1); }",
            "        .slide img { width: 100%; max-width: 640px; border: 1px solid #ddd; border-radius: 4px; }",
            "        .slide-footer { margin-top: 8px; color: #555; font-size: 14px; }",
            "    </style>",
            "</head>",
            "<body>",
            "    <h1>改進模式幻燈片預覽</h1>",
            "    <div class=\"summary\">",
            f"        <p><strong>影片：</strong>{os.path.basename(metadata.get('video_path', '未知來源'))}</p>",
            f"        <p><strong>總幻燈片數：</strong>{len(metadata.get('slides', []))}</p>",
            f"        <p><strong>影片時長 (推估)：</strong>{duration:.1f} 秒</p>",
            "    </div>",
        ]

        for slide in metadata.get("slides", []):
            filename = slide.get("filename", "")
            timestamp = slide.get("timestamp", 0.0)
            group_id = slide.get("group_id", -1)
            html_lines.extend([
                "    <div class=\"slide\">",
                f"        <img src=\"{filename}\" alt=\"{filename}\">",
                "        <div class=\"slide-footer\">",
                f"            <p>時間戳：{timestamp:.2f} 秒 | 分組：{group_id if group_id is not None else -1}</p>",
                f"            <p>檔案：{filename}</p>",
                "        </div>",
                "    </div>",
            ])

        html_lines.extend(["</body>", "</html>"])
        html_path = os.path.join(slides_folder, "preview.html")
        with open(html_path, "w", encoding="utf-8") as handle:
            handle.write("\n".join(html_lines))

    def process_batch(self, path: str):
        """批次處理指定路徑"""
        print("\n🎬 改進模式批次幻燈片捕獲工具")
        print("=" * 60)

        print("\n搜尋影片中...")
        video_files = self.find_videos(path)
        if not video_files:
            print("未找到符合條件的影片")
            return

        self.stats["total_videos"] = len(video_files)

        videos_to_process: List[str] = []
        videos_to_skip: List[Dict[str, str]] = []

        for video in video_files:
            video_dir = os.path.dirname(video)
            video_name = os.path.splitext(os.path.basename(video))[0]
            if not self.force:
                skip = False
                for folder_name in [
                    f"slides_improved_{video_name}",
                    f"{video_name}_slides_improved",
                    f"{video_name}_slides",
                    f"{video_name}_slide",
                    "slides_improved",
                    "slides",
                    "slide",
                    "Slides",
                    "Slide",
                ]:
                    folder_path = os.path.join(video_dir, folder_name)
                    if os.path.isdir(folder_path):
                        try:
                            files = os.listdir(folder_path)
                            if any(
                                f.lower().endswith((".jpg", ".jpeg", ".png"))
                                for f in files
                                if not f.startswith("._")
                            ):
                                videos_to_skip.append({
                                    "video": video,
                                    "reason": "slides",
                                    "detail": folder_name,
                                })
                                skip = True
                                break
                        except OSError:
                            pass
                if skip:
                    continue

            videos_to_process.append(video)

        print(f"\n找到 {len(video_files)} 支影片，其中 {len(videos_to_process)} 支需要處理")
        if videos_to_skip:
            print(f"{len(videos_to_skip)} 支將被跳過 (已有幻燈片)：")
            for item in videos_to_skip:
                base = os.path.basename(item["video"])
                print(f"  - {base} -> {item['detail']} (已有幻燈片)")

        if self.list_only:
            print("\n僅列表模式，不執行實際處理")
            return

        if videos_to_process and len(videos_to_process) > 1 and not self.yes:
            confirm = input(f"\n確定要處理 {len(videos_to_process)} 支影片嗎？(y/n): ")
            if confirm.lower() != "y":
                print("已取消批次處理")
                return
        elif not videos_to_process:
            print("\n沒有需要處理的影片")
            return

        batch_start = time.time()
        for index, video_path in enumerate(videos_to_process, start=1):
            print(f"\n進度: {index}/{len(videos_to_process)}")
            result = self.process_video(video_path)
            if result["status"] == "success":
                self.processed_videos.append(result)
                self.stats["processed"] += 1
            elif result["status"] == "failed":
                self.failed_videos.append(result)
                self.stats["failed"] += 1

        self.stats["processing_time"] = time.time() - batch_start
        self.show_summary()

    def show_summary(self):
        print("\n" + "=" * 60)
        print("📊 批次處理總結")
        print("=" * 60)

        skipped = self.stats["total_videos"] - self.stats["processed"] - self.stats["failed"]
        print(f"  總影片數: {self.stats['total_videos']}")
        print(f"  成功: {self.stats['processed']}")
        print(f"  失敗: {self.stats['failed']}")
        print(f"  跳過: {skipped}")
        print(f"  總幻燈片數: {self.stats['total_slides']}")
        print(f"  總推估分組數: {self.stats['total_groups']}")
        print(f"  總耗時: {self.stats['processing_time']:.1f} 秒")

        if self.stats["processed"] > 0:
            avg_time = self.stats["processing_time"] / self.stats["processed"]
            avg_slides = self.stats["total_slides"] / self.stats["processed"]
            print(f"  平均每支影片耗時: {avg_time:.1f} 秒")
            print(f"  平均每支影片幻燈片: {avg_slides:.1f}")

        if self.processed_videos:
            print("\n✅ 成功處理：")
            for result in self.processed_videos:
                print(f"  - {os.path.basename(result['video'])}")
                print(f"    幻燈片: {result['slide_count']} | 推估分組: {result['group_count']}")
                print(f"    輸出: {result['output']}")
                if result.get("ppt_file"):
                    print(f"    PowerPoint: {result['ppt_file']}")

        if self.failed_videos:
            print("\n❌ 處理失敗：")
            for result in self.failed_videos:
                print(f"  - {os.path.basename(result['video'])}")
                print(f"    錯誤: {result['error']}")


def main():
    parser = argparse.ArgumentParser(
        description="批量使用改進模式截取投影片",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""範例：
  # 處理單支影片
  %(prog)s /path/to/video.mp4

  # 處理整個資料夾 (非遞迴)
  %(prog)s /path/to/folder

  # 遞迴搜尋子資料夾
  %(prog)s /path/to/folder --recursive

  # 自訂閾值並自動挑選代表幻燈片
  %(prog)s /path/to/folder --threshold 0.88 --auto-select
""",
    )

    parser.add_argument("path", help="影片檔或資料夾路徑")
    parser.add_argument(
        "-r",
        "--recursive",
        action="store_true",
        help="遞迴搜尋子資料夾中的影片",
    )
    parser.add_argument(
        "-t",
        "--threshold",
        type=float,
        default=0.85,
        help="相似度閾值 (預設: 0.85)",
    )
    parser.add_argument(
        "-a",
        "--auto-select",
        action="store_true",
        help="自動挑選代表幻燈片至 selected_slides/",
    )
    parser.add_argument(
        "-f",
        "--force",
        action="store_true",
        help="強制重新處理 (忽略既有輸出)",
    )
    parser.add_argument(
        "-l",
        "--list-only",
        action="store_true",
        help="僅列出將處理的影片，不執行實際處理",
    )
    parser.add_argument(
        "-y",
        "--yes",
        action="store_true",
        help="自動確認批次處理，不顯示提示",
    )
    parser.add_argument(
        "--no-ppt",
        dest="generate_ppt",
        action="store_false",
        help="不自動將幻燈片轉換為 PowerPoint",
    )
    parser.set_defaults(generate_ppt=True)

    args = parser.parse_args()

    processor = BatchImprovedSlideCapture(
        threshold=args.threshold,
        auto_select=args.auto_select,
        recursive=args.recursive,
        force=args.force,
        list_only=args.list_only,
        yes=args.yes,
        generate_ppt=args.generate_ppt,
    )

    try:
        processor.process_batch(args.path)
    except KeyboardInterrupt:
        print("\n\n⚠️  使用者中斷處理")
        processor.show_summary()
    except Exception as exc:
        print(f"\n❌ 發生錯誤: {exc}")
        import traceback

        traceback.print_exc()


if __name__ == "__main__":
    main()
