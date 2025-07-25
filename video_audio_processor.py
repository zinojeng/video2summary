#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
視頻和音頻處理工具

此程序可以處理以下功能：
1. 從 .mp4 文件中提取音頻
2. 捕獲視頻中的幻燈片，並生成 PowerPoint 或 Markdown 文件
"""

import os
import sys
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
import subprocess
import traceback
import threading
from PIL import Image, ImageTk
import importlib
import json
from collections import defaultdict

# 需要時會動態導入的模塊:
# - moviepy (extract_audio_from_video)
# - cv2, numpy, skimage.metrics (capture_slides_from_video)
# - pptx (generate_ppt_from_images)
# - markitdown_helper (custom module for image to markdown conversion)
# - importlib.util (for checking installed modules)

def check_dependencies():
    """檢查必要依賴是否已安裝"""
    import importlib.util
    
    required_packages = [
        "opencv-python", "numpy", "pillow", "moviepy",
        "python-pptx", "scikit-image"
    ]
    
    # 分開檢查 markitdown，因為它可選
    optional_packages = ["markitdown"]
    
    missing_packages = []
    missing_optional = []
    
    for package in required_packages:
        try:
            # 處理特殊套件名稱
            import_name = package.replace("-", "_")
            # 處理 opencv-python 特例
            if package == "opencv-python":
                import_name = "cv2"
            # 處理 python-pptx 特例
            elif package == "python-pptx":
                import_name = "pptx"
            # 處理 pillow 特例
            elif package == "pillow":
                import_name = "PIL"
            # 處理 scikit-image 特例
            elif package == "scikit-image":
                import_name = "skimage"
                
            spec = importlib.util.find_spec(import_name)
            if spec is None:
                missing_packages.append(package)
        except ImportError:
            missing_packages.append(package)
    
    for package in optional_packages:
        try:
            spec = importlib.util.find_spec(package)
            if spec is None:
                missing_optional.append(package)
        except ImportError:
            missing_optional.append(package)
    
    if missing_optional:
        print(f"注意: 未安裝可選套件: {', '.join(missing_optional)}")
        print("如果您想使用 MarkItDown 處理功能，請安裝:")
        print("pip install markitdown>=0.1.1")
    
    return missing_packages


def install_dependencies(packages):
    """安裝缺失的依賴"""
    print(f"正在安裝必要依賴: {', '.join(packages)}")
    
    try:
        subprocess.check_call(
            [sys.executable, "-m", "pip", "install"] + packages
        )
        return True
    except subprocess.CalledProcessError:
        return False


def extract_audio_from_video(video_path, output_path=None, format="mp3"):
    """
    從視頻文件中提取音頻
    
    參數:
        video_path: 視頻文件路徑
        output_path: 輸出音頻文件路徑，默認為原視頻名稱加上音頻格式後綴
        format: 輸出音頻格式，默認為 mp3
        
    返回:
        success: 是否成功
        output_file: 輸出文件路徑或錯誤信息
    """
    try:
        from moviepy import VideoFileClip
        
        # 如果未指定輸出路徑，使用默認路徑
        if not output_path:
            base_name = os.path.splitext(video_path)[0]
            output_path = f"{base_name}.{format}"
            
        # 打開視頻文件
        video_clip = VideoFileClip(video_path)
        
        # 提取音頻
        audio_clip = video_clip.audio
        
        # 寫入音頻文件
        audio_clip.write_audiofile(output_path)
        
        # 釋放資源
        audio_clip.close()
        video_clip.close()
        
        return True, output_path
        
    except Exception as e:
        error_msg = str(e)
        print(f"提取音頻時出錯: {error_msg}")
        return False, error_msg


def calculate_phash(img, hash_size=8):
    """計算感知哈希（pHash）"""
    import cv2
    import numpy as np
    
    # 轉換為灰度圖
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY) if len(img.shape) == 3 else img
    
    # 調整大小到 hash_size x hash_size
    resized = cv2.resize(gray, (hash_size, hash_size))
    
    # 計算DCT
    dct_result = cv2.dct(np.float32(resized))
    
    # 只保留左上角的低頻部分
    dct_low = dct_result[:hash_size, :hash_size]
    
    # 計算平均值（排除第一個元素）
    avg = np.mean(dct_low[1:, 1:])
    
    # 生成哈希
    hash_bits = (dct_low > avg).flatten()
    
    # 轉換為十六進制字符串
    hash_int = 0
    for bit in hash_bits:
        hash_int = (hash_int << 1) | int(bit)
    
    return format(hash_int, f'0{hash_size*hash_size//4}x')


def capture_slides_from_video(video_path, output_folder=None, similarity_threshold=0.8, enable_metadata=True):
    """
    從視頻文件中捕獲幻燈片
    
    參數:
        video_path: 視頻文件路徑
        output_folder: 輸出幻燈片圖片的文件夾，默認為 video_slides
        similarity_threshold: 圖像相似度閾值，用於檢測幻燈片變化，值越低檢測越敏感
        enable_metadata: 是否生成元數據文件
        
    返回:
        success: 是否成功
        result: 包含提取信息的字典或錯誤信息
    """
    try:
        import cv2
        import numpy as np
        from skimage.metrics import structural_similarity as ssim
        
        # 如果未指定輸出文件夾，使用默認文件夾
        if not output_folder:
            # 在視頻所在目錄創建slides文件夾
            video_dir = os.path.dirname(video_path)
            output_folder = os.path.join(video_dir, "slides")
        
        # 確保輸出文件夾存在
        os.makedirs(output_folder, exist_ok=True)
            
        # 打開視頻文件
        cap = cv2.VideoCapture(video_path)
        
        if not cap.isOpened():
            return False, "無法打開視頻文件"
            
        # 獲取視頻信息
        fps = cap.get(cv2.CAP_PROP_FPS)
        frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        duration = frame_count / fps
        
        # 計算採樣間隔（秒）
        # 每秒檢查一次幀，可以根據需要調整
        sample_interval = 1
        
        # 初始化幀計數器和上一幀
        frame_idx = 0
        prev_frame = None
        saved_count = 0
        last_saved_time = -1000  # 上次保存的時間點，初始為負值
        
        # 元數據相關
        slides_data = []
        frame_hashes = {}  # 存儲幀哈希

        while True:
            # 設置當前幀位置
            cap.set(cv2.CAP_PROP_POS_FRAMES, frame_idx)
            
            # 讀取當前幀
            ret, frame = cap.read()
            
            if not ret:
                break
                
            # 計算當前時間點（秒）
            current_time = frame_idx / fps
            
            # 轉換為灰度圖像用於相似度比較
            gray_frame = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
            
            # 如果是第一幀或與上一幀相比相似度低於閾值，則保存
            if prev_frame is None or (
                current_time - last_saved_time >= 2.0 and  # 確保至少間隔2秒
                ssim(prev_frame, gray_frame) < similarity_threshold
            ):
                # 計算感知哈希
                phash = calculate_phash(frame)
                
                # 檢查是否為重複幀
                is_duplicate = False
                for stored_hash in frame_hashes.values():
                    # 計算哈希相似度
                    hamming_dist = bin(int(phash, 16) ^ int(stored_hash, 16)).count('1')
                    if hamming_dist < 5:  # 相似度闾值
                        is_duplicate = True
                        break
                
                if not is_duplicate:
                    # 轉換時間格式
                    minutes = int(current_time / 60)
                    seconds = current_time % 60
                    
                    # 生成文件名
                    filename = f"slide_{saved_count:03d}_t{minutes}m{seconds:.1f}s_h{phash[:8]}.jpg"
                    output_path = os.path.join(output_folder, filename)
                    
                    # 保存圖片
                    cv2.imwrite(output_path, frame, [cv2.IMWRITE_JPEG_QUALITY, 95])
                    
                    # 記錄元數據
                    frame_hashes[frame_idx] = phash
                    slides_data.append({
                        'index': saved_count,
                        'filename': filename,
                        'frame_index': frame_idx,
                        'timestamp': current_time,
                        'phash': phash
                    })
                    
                    saved_count += 1
                    last_saved_time = current_time
                    
                    # 更新上一幀
                    prev_frame = gray_frame
                    
                    print(f"保存幻燈片 {saved_count}: {filename}")
                else:
                    print(f"跳過重複幀 於 {current_time:.2f} 秒")
            
            # 更新幀索引，跳過一些幀以提高效率
            frame_idx += int(fps * sample_interval)
            
            # 檢查是否已到達視頻結尾
            if frame_idx >= frame_count:
                break
                
        # 釋放資源
        cap.release()
        
        # 保存元數據
        if enable_metadata and saved_count > 0:
            metadata_path = os.path.join(output_folder, 'slides_metadata.json')
            with open(metadata_path, 'w', encoding='utf-8') as f:
                json.dump({
                    'video_path': video_path,
                    'total_frames': frame_count,
                    'fps': fps,
                    'threshold': similarity_threshold,
                    'slides': slides_data
                }, f, indent=2, ensure_ascii=False)
            print(f"\n元數據已保存到: {metadata_path}")
        
        return True, {
            "output_folder": output_folder,
            "slide_count": saved_count,
            "video_duration": duration
        }
        
    except Exception as e:
        error_msg = str(e)
        traceback.print_exc()
        return False, error_msg


def generate_ppt_from_images(image_folder, output_file=None, title="視頻捕獲的幻燈片"):
    """
    將圖片文件夾轉換為 PowerPoint 文件
    
    參數:
        image_folder: 包含幻燈片圖片的文件夾
        output_file: 輸出的 PowerPoint 文件路徑，默認為文件夾名+.pptx
        title: 演示文稿標題
        
    返回:
        success: 是否成功
        output_file: 輸出文件路徑或錯誤信息
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
        
        # 如果未指定輸出文件，使用默認文件
        if not output_file:
            output_file = os.path.join(
                os.path.dirname(image_folder), 
                f"{os.path.basename(image_folder)}.pptx"
            )
            
        # 創建新的 PowerPoint 演示文稿
        prs = Presentation()
        
        # 設置幻燈片尺寸為 16:9
        prs.slide_width = Inches(10)  # 16:9 寬度
        prs.slide_height = Inches(5.625)  # 16:9 高度
        
        # 獲取幻燈片尺寸
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # 獲取所有圖片文件（過濾掉 macOS 隱藏文件）
        image_files = []
        for filename in sorted(os.listdir(image_folder)):
            # 跳過 macOS 隱藏文件
            if filename.startswith('._'):
                continue
            if filename.lower().endswith(('.png', '.jpg', '.jpeg')):
                image_files.append(os.path.join(image_folder, filename))
                
        if not image_files:
            return False, "未找到圖片文件"
            
        # 為每張圖片創建一個幻燈片（不添加標題幻燈片）
        blank_slide_layout = prs.slide_layouts[6]  # 空白布局
        
        for img_path in image_files:
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 添加圖片
            left = top = Inches(0)
            img = slide.shapes.add_picture(
                img_path, left, top, width=slide_width, height=slide_height
            )
            
        # 保存演示文稿
        prs.save(output_file)
        
        return True, output_file
        
    except Exception as e:
        error_msg = str(e)
        print(f"生成 PowerPoint 時出錯: {error_msg}")
        traceback.print_exc()
        return False, error_msg


def generate_markdown_from_images(image_folder, output_file=None, title="視頻捕獲的幻燈片", use_markitdown=True, api_key=None):
    """
    將圖片文件夾轉換為 Markdown 文件
    
    參數:
        image_folder: 包含幻燈片圖片的文件夾
        output_file: 輸出的 Markdown 文件路徑，默認為文件夾名+.md
        title: Markdown 文件標題
        use_markitdown: 是否使用 MarkItDown 庫進行圖片文本提取
        api_key: 如果使用 MarkItDown 並需要 LLM 支持，提供 API Key
        
    返回:
        success: 是否成功
        output_file: 輸出文件路徑或錯誤信息
    """
    try:
        # 如果未指定輸出文件，使用默認文件
        if not output_file:
            output_file = os.path.join(
                os.path.dirname(image_folder), 
                f"{os.path.basename(image_folder)}.md"
            )
            
        # 獲取所有圖片文件（過濾掉 macOS 隱藏文件）
        image_files = []
        for filename in sorted(os.listdir(image_folder)):
            # 跳過 macOS 隱藏文件
            if filename.startswith('._'):
                continue
            if filename.lower().endswith(('.png', '.jpg', '.jpeg')):
                image_files.append(os.path.join(image_folder, filename))
                
        if not image_files:
            return False, "未找到圖片文件"
        
        # 嘗試使用我們的自定義 markitdown_helper
        try:
            import markitdown_helper
            
            print(f"使用自定義 markitdown_helper 處理 {len(image_files)} 張圖片...")
            success, result, info = markitdown_helper.convert_images_to_markdown(
                image_paths=image_files,
                output_file=output_file,
                title=title,
                use_llm=(api_key is not None),
                api_key=api_key
            )
            
            if success:
                return True, result
            else:
                print(f"使用 markitdown_helper 失敗: {info.get('error', '未知錯誤')}")
                # 繼續使用基本方法
        except ImportError:
            print("找不到 markitdown_helper 模組，繼續嘗試其他方法...")
        except Exception as e:
            print(f"使用 markitdown_helper 時出錯: {e}")
            traceback.print_exc()
            
        # 如果使用 MarkItDown，則調用相應函數
        if use_markitdown:
            try:
                # 檢查 markitdown 是否已安裝
                import importlib.util
                markitdown_spec = importlib.util.find_spec("markitdown")
                
                if markitdown_spec is not None:
                    import markitdown
                    
                    # 使用 markitdown 庫
                    print(f"使用 MarkItDown 庫處理 {len(image_files)} 張圖片...")
                    
                    # 創建 MarkItDown 實例
                    converter = markitdown.MarkItDown()
                    
                    # 添加自訂標題的 Markdown 內容
                    with open(output_file, 'w', encoding='utf-8') as f:
                        f.write(f"# {title}\n\n")
                        
                        for i, img_path in enumerate(image_files):
                            # 獲取相對路徑
                            try:
                                rel_path = os.path.relpath(
                                    img_path,
                                    os.path.dirname(output_file)
                                )
                            except Exception:
                                rel_path = img_path
                                
                            # 添加幻燈片標題和圖片
                            slide_num = i + 1
                            f.write(f"## 幻燈片 {slide_num}\n\n")
                            f.write(f"![幻燈片 {slide_num}]({rel_path})\n\n")
                            
                            try:
                                # 使用 markitdown 分析圖片內容
                                result = converter.convert(img_path)
                                if result and result.text_content:
                                    f.write(f"{result.text_content.strip()}\n\n")
                            except Exception as img_err:
                                print(f"分析圖片 {img_path} 時出錯: {img_err}")
                            
                            f.write("---\n\n")
                    
                    print(f"已成功生成 Markdown 文件: {output_file}")
                    return True, output_file
                else:
                    print("找不到 markitdown 模組，將使用基本方法生成 Markdown")
                    use_markitdown = False
                    
            except ImportError as e:
                print(f"MarkItDown 庫導入錯誤: {e}")
                print("將使用基本方法生成 Markdown")
                use_markitdown = False
            except Exception as e:
                print(f"使用 MarkItDown 時出錯: {e}")
                print("將使用基本方法生成 Markdown")
                traceback.print_exc()
                use_markitdown = False
        
        # 使用基本方法生成 Markdown
        print("使用基本方法生成 Markdown...")
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(f"# {title}\n\n")
            
            for i, img_path in enumerate(image_files):
                # 獲取相對路徑，使 Markdown 中的圖片鏈接更短
                try:
                    rel_path = os.path.relpath(
                        img_path, 
                        os.path.dirname(output_file)
                    )
                except Exception:
                    rel_path = img_path
                    
                # 添加幻燈片標題和圖片
                slide_num = i + 1
                f.write(f"## 幻燈片 {slide_num}\n\n")
                f.write(f"![幻燈片 {slide_num}]({rel_path})\n\n")
                f.write("---\n\n")
        
        return True, output_file
        
    except Exception as e:
        error_msg = str(e)
        print(f"生成 Markdown 時出錯: {error_msg}")
        traceback.print_exc()
        return False, error_msg


class VideoAudioProcessor:
    """視頻和音頻處理應用"""
    
    def __init__(self, root):
        self.root = root
        self.root.title("視頻和音頻處理工具")
        self.root.geometry("900x700")
        
        # 保存使用者的 API Key
        self.saved_api_key = os.environ.get("OPENAI_API_KEY", "")
        
        # 創建頁面框架
        self.setup_ui()
    
    def setup_ui(self):
        """設置使用者介面"""
        # 建立選項卡
        self.notebook = ttk.Notebook(self.root)
        self.notebook.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        
        # 第一個選項卡：音頻提取
        self.audio_frame = ttk.Frame(self.notebook)
        self.notebook.add(self.audio_frame, text="音頻提取")
        
        # 第二個選項卡：幻燈片捕獲
        self.slide_frame = ttk.Frame(self.notebook)
        self.notebook.add(self.slide_frame, text="幻燈片捕獲")
        
        # 第三個選項卡：幻燈片處理
        self.process_frame = ttk.Frame(self.notebook)
        self.notebook.add(self.process_frame, text="幻燈片處理")
        
        # 設置音頻提取 UI
        self.setup_audio_ui()
        
        # 設置幻燈片捕獲 UI
        self.setup_slide_ui()
        
        # 設置幻燈片處理 UI
        self.setup_process_ui()
    
    def setup_audio_ui(self):
        """設置音頻提取界面"""
        frame = self.audio_frame
        
        # 頂部說明文字
        info_label = tk.Label(
            frame, 
            text="此功能可以從視頻文件中提取音頻軌道",
            font=("Arial", 12)
        )
        info_label.pack(pady=20)
        
        # 視頻文件選擇區域
        video_frame = tk.Frame(frame)
        video_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(video_frame, text="視頻文件:").pack(side=tk.LEFT, padx=10)
        self.video_entry = tk.Entry(video_frame, width=50)
        self.video_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        
        self.video_browse_btn = tk.Button(
            video_frame, text="瀏覽...", 
            command=lambda: self.browse_file(
                self.video_entry, 
                filetypes=[("視頻文件", "*.mp4 *.avi *.mkv *.mov")]
            )
        )
        self.video_browse_btn.pack(side=tk.LEFT, padx=5)
        
        # 輸出音頻設置區域
        output_frame = tk.Frame(frame)
        output_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(output_frame, text="輸出文件:").pack(side=tk.LEFT, padx=10)
        self.audio_output_entry = tk.Entry(output_frame, width=50)
        self.audio_output_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        
        self.audio_browse_btn = tk.Button(
            output_frame, text="瀏覽...", 
            command=lambda: self.save_file(
                self.audio_output_entry, 
                filetypes=[("音頻文件", "*.mp3 *.wav *.aac")]
            )
        )
        self.audio_browse_btn.pack(side=tk.LEFT, padx=5)
        
        # 音頻格式選擇
        format_frame = tk.Frame(frame)
        format_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(format_frame, text="音頻格式:").pack(side=tk.LEFT, padx=10)
        
        self.audio_format_var = tk.StringVar(value="mp3")
        
        tk.Radiobutton(
            format_frame, text="MP3", 
            variable=self.audio_format_var, value="mp3"
        ).pack(side=tk.LEFT, padx=5)
        
        tk.Radiobutton(
            format_frame, text="WAV", 
            variable=self.audio_format_var, value="wav"
        ).pack(side=tk.LEFT, padx=5)
        
        tk.Radiobutton(
            format_frame, text="AAC", 
            variable=self.audio_format_var, value="aac"
        ).pack(side=tk.LEFT, padx=5)
        
        # 狀態顯示
        self.audio_status_var = tk.StringVar(value="準備就緒")
        status_label = tk.Label(
            frame, textvariable=self.audio_status_var,
            font=("Arial", 10), fg="blue"
        )
        status_label.pack(pady=10)
        
        # 進度條
        self.audio_progress = ttk.Progressbar(
            frame, orient="horizontal", length=300, mode="indeterminate"
        )
        self.audio_progress.pack(pady=10)
        
        # 提取按鈕
        self.extract_btn = tk.Button(
            frame, text="提取音頻", command=self.extract_audio,
            bg="#4CAF50", fg="white", height=2, width=20
        )
        self.extract_btn.pack(pady=20)
        
        # 分隔線
        separator = ttk.Separator(frame, orient='horizontal')
        separator.pack(fill=tk.X, pady=20)
        
        # 音頻轉文字說明
        transcribe_label = tk.Label(
            frame, 
            text="音頻轉文字功能（需要 OpenAI API Key）",
            font=("Arial", 12, "bold")
        )
        transcribe_label.pack(pady=10)
        
        # 音頻文件選擇區域
        audio_file_frame = tk.Frame(frame)
        audio_file_frame.pack(fill=tk.X, pady=5)
        
        tk.Label(audio_file_frame, text="音頻文件:").pack(side=tk.LEFT, padx=10)
        self.transcribe_audio_entry = tk.Entry(audio_file_frame, width=50)
        self.transcribe_audio_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        
        self.transcribe_browse_btn = tk.Button(
            audio_file_frame, text="瀏覽...", 
            command=lambda: self.browse_file(
                self.transcribe_audio_entry, 
                filetypes=[("音頻文件", "*.mp3 *.mp4 *.mpeg *.mpga *.m4a *.wav *.webm")]
            )
        )
        self.transcribe_browse_btn.pack(side=tk.LEFT, padx=5)
        
        # 模型選擇
        model_frame = tk.Frame(frame)
        model_frame.pack(fill=tk.X, pady=5)
        
        tk.Label(model_frame, text="模型選擇:").pack(side=tk.LEFT, padx=10)
        
        self.transcribe_model_var = tk.StringVar(value="gpt-4o-transcribe")
        
        tk.Radiobutton(
            model_frame, text="GPT-4o", 
            variable=self.transcribe_model_var, value="gpt-4o-transcribe"
        ).pack(side=tk.LEFT, padx=5)
        
        tk.Radiobutton(
            model_frame, text="GPT-4o-Mini", 
            variable=self.transcribe_model_var, value="gpt-4o-mini-transcribe"
        ).pack(side=tk.LEFT, padx=5)
        
        # 輸出格式選擇
        format_frame2 = tk.Frame(frame)
        format_frame2.pack(fill=tk.X, pady=5)
        
        tk.Label(format_frame2, text="輸出格式:").pack(side=tk.LEFT, padx=10)
        
        self.transcribe_format_var = tk.StringVar(value="text")
        
        tk.Radiobutton(
            format_frame2, text="純文字 (.txt)", 
            variable=self.transcribe_format_var, value="text"
        ).pack(side=tk.LEFT, padx=5)
        
        tk.Radiobutton(
            format_frame2, text="字幕 (.srt)", 
            variable=self.transcribe_format_var, value="srt"
        ).pack(side=tk.LEFT, padx=5)
        
        tk.Radiobutton(
            format_frame2, text="Markdown (.md)", 
            variable=self.transcribe_format_var, value="markdown"
        ).pack(side=tk.LEFT, padx=5)
        
        # 轉錄狀態顯示
        self.transcribe_status_var = tk.StringVar(value="準備就緒")
        transcribe_status_label = tk.Label(
            frame, textvariable=self.transcribe_status_var,
            font=("Arial", 10), fg="blue"
        )
        transcribe_status_label.pack(pady=5)
        
        # 轉錄進度條
        self.transcribe_progress = ttk.Progressbar(
            frame, orient="horizontal", length=300, mode="indeterminate"
        )
        self.transcribe_progress.pack(pady=5)
        
        # 轉錄按鈕
        self.transcribe_btn = tk.Button(
            frame, text="開始轉錄", command=self.transcribe_audio,
            bg="#2196F3", fg="white", height=2, width=20
        )
        self.transcribe_btn.pack(pady=10)
    
    def browse_file(self, entry_widget, filetypes):
        """瀏覽並選擇文件"""
        file_path = filedialog.askopenfilename(
            initialdir=".",
            title="選擇文件",
            filetypes=filetypes
        )
        if file_path:
            entry_widget.delete(0, tk.END)
            entry_widget.insert(0, file_path)
            
            # 如果是視頻文件，自動生成音頻輸出路徑
            if entry_widget == self.video_entry:
                base_name = os.path.splitext(file_path)[0]
                audio_format = self.audio_format_var.get()
                output_path = f"{base_name}.{audio_format}"
                
                self.audio_output_entry.delete(0, tk.END)
                self.audio_output_entry.insert(0, output_path)
            
            # 如果是幻燈片捕獲的視頻文件，自動設置輸出文件夾
            elif entry_widget == self.slide_video_entry:
                # 獲取視頻文件所在的目錄
                video_dir = os.path.dirname(file_path)
                # 在視頻所在目錄創建輸出文件夾
                output_folder = os.path.join(video_dir, "slides")
                
                self.slide_output_entry.delete(0, tk.END)
                self.slide_output_entry.insert(0, output_folder)
    
    def save_file(self, entry_widget, filetypes):
        """選擇保存文件的路徑"""
        # 獲取當前輸入框中的路徑
        current_path = entry_widget.get()
        
        # 確定初始目錄
        if current_path:
            initial_dir = os.path.dirname(current_path)
            if not os.path.exists(initial_dir):
                initial_dir = "."
        else:
            initial_dir = "."
            
        file_path = filedialog.asksaveasfilename(
            initialdir=initial_dir,
            title="保存文件",
            filetypes=filetypes
        )
        if file_path:
            entry_widget.delete(0, tk.END)
            entry_widget.insert(0, file_path)
    
    def extract_audio(self):
        """提取音頻按鈕點擊處理"""
        video_path = self.video_entry.get()
        output_path = self.audio_output_entry.get()
        audio_format = self.audio_format_var.get()
        
        if not video_path:
            messagebox.showwarning("警告", "請選擇視頻文件")
            return
            
        # 確保輸出路徑有正確的擴展名
        if output_path and not output_path.lower().endswith(f".{audio_format}"):
            output_path = f"{os.path.splitext(output_path)[0]}.{audio_format}"
            self.audio_output_entry.delete(0, tk.END)
            self.audio_output_entry.insert(0, output_path)
        
        # 開始提取（在背景線程中運行）
        self.audio_status_var.set("正在提取音頻...")
        self.audio_progress.start(10)
        self.extract_btn.config(state=tk.DISABLED)
        
        def extract_thread():
            success, result = extract_audio_from_video(
                video_path, output_path, audio_format
            )
            
            # 在主線程中更新 UI
            self.root.after(0, lambda: self.extraction_completed(success, result))
        
        threading.Thread(target=extract_thread).start()
    
    def extraction_completed(self, success, result):
        """音頻提取完成後的處理"""
        self.audio_progress.stop()
        self.extract_btn.config(state=tk.NORMAL)
        
        if success:
            self.audio_status_var.set(f"提取完成: {result}")
            messagebox.showinfo("成功", f"音頻已成功提取到: {result}")
        else:
            self.audio_status_var.set(f"提取失敗: {result}")
            messagebox.showerror("錯誤", f"音頻提取失敗: {result}")
    
    def transcribe_audio(self):
        """開始音頻轉錄"""
        audio_path = self.transcribe_audio_entry.get()
        
        if not audio_path:
            messagebox.showwarning("警告", "請選擇音頻文件")
            return
        
        if not os.path.exists(audio_path):
            messagebox.showerror("錯誤", "音頻文件不存在")
            return
        
        # 檢查 API Key
        api_key = self.saved_api_key or os.environ.get("OPENAI_API_KEY", "")
        if not api_key:
            # 彈出對話框要求輸入 API Key
            dialog = tk.Toplevel(self.root)
            dialog.title("輸入 OpenAI API Key")
            dialog.geometry("500x150")
            
            tk.Label(dialog, text="請輸入您的 OpenAI API Key:").pack(pady=10)
            api_key_entry = tk.Entry(dialog, width=60, show="*")
            api_key_entry.pack(pady=10)
            
            def confirm_api_key():
                key = api_key_entry.get().strip()
                if key:
                    self.saved_api_key = key
                    dialog.destroy()
                    self.start_transcription(audio_path)
                else:
                    messagebox.showwarning("警告", "請輸入有效的 API Key")
            
            tk.Button(dialog, text="確認", command=confirm_api_key).pack(pady=10)
            
            dialog.transient(self.root)
            dialog.grab_set()
            self.root.wait_window(dialog)
        else:
            self.start_transcription(audio_path)
    
    def start_transcription(self, audio_path):
        """開始轉錄處理"""
        model = self.transcribe_model_var.get()
        output_format = self.transcribe_format_var.get()
        
        # 生成輸出文件名
        base_name = os.path.splitext(audio_path)[0]
        if output_format == "text":
            output_path = f"{base_name}_transcription.txt"
        elif output_format == "srt":
            output_path = f"{base_name}_transcription.srt"
        else:  # markdown
            output_path = f"{base_name}_transcription.md"
        
        # 更新狀態
        self.transcribe_status_var.set("正在轉錄音頻...")
        self.transcribe_progress.start(10)
        self.transcribe_btn.config(state=tk.DISABLED)
        
        def transcribe_thread():
            try:
                # 調用轉錄函數
                result = self.transcribe_audio_to_text(
                    audio_path, 
                    self.saved_api_key,
                    model,
                    output_format
                )
                
                # 保存結果到文件
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(result)
                
                # 在主線程中更新 UI
                self.root.after(0, lambda: self.transcription_completed(True, output_path))
                
            except Exception as e:
                error_msg = str(e)
                self.root.after(0, lambda: self.transcription_completed(False, error_msg))
        
        threading.Thread(target=transcribe_thread).start()
    
    def transcribe_audio_to_text(self, file_path, api_key, model="gpt-4o-transcribe", output_format="text"):
        """
        使用 GPT-4o 模型轉錄音頻（改進版）
        
        Args:
            file_path: 音頻檔案路徑
            api_key: OpenAI API 金鑰
            model: 模型名稱 (gpt-4o-transcribe 或 gpt-4o-mini-transcribe)
            output_format: 輸出格式 ('text', 'srt', 'markdown')
        
        Returns:
            轉錄結果文字
        """
        try:
            # 首先嘗試使用改進的轉錄模組
            try:
                from gpt4o_transcribe_improved import AudioTranscriber
                
                # 使用改進的轉錄器
                transcriber = AudioTranscriber(api_key)
                result = transcriber.transcribe(
                    file_path,
                    model=model,
                    language="zh",
                    output_format=output_format,
                    auto_convert=True  # 自動轉換格式以提高相容性
                )
                return result
                
            except ImportError:
                # 如果改進模組不存在，使用原始方法
                from openai import OpenAI
                
                client = OpenAI(api_key=api_key)
                
                # 根據格式設定 response_format
                response_format = "text"
                if output_format == "srt":
                    response_format = "srt"
                elif output_format == "markdown":
                    response_format = "text"  # 先取得文字再轉換為 markdown
                
                with open(file_path, "rb") as audio_file:
                    transcript = client.audio.transcriptions.create(
                        model=model,
                        file=audio_file,
                        response_format=response_format
                    )
                
                if output_format == "markdown":
                    # 將文字轉換為 markdown 格式
                    if hasattr(transcript, 'text'):
                        result = f"# 語音轉錄結果\n\n{transcript.text}\n"
                    else:
                        result = f"# 語音轉錄結果\n\n{transcript}\n"
                elif output_format == "srt":
                    # SRT 格式直接回傳字串
                    if hasattr(transcript, 'text'):
                        result = transcript.text
                    else:
                        result = transcript
                else:
                    # 純文字格式
                    if hasattr(transcript, 'text'):
                        result = transcript.text
                    else:
                        result = transcript
                    
                return result
            
        except ImportError:
            raise Exception("請先安裝 OpenAI 套件: pip install openai>=1.0.0")
        except Exception as e:
            # 如果是檔案格式錯誤，提供更詳細的說明
            if "corrupted" in str(e) or "unsupported" in str(e):
                raise Exception(
                    f"GPT-4o 轉錄失敗: {str(e)}\n\n"
                    "可能的解決方法:\n"
                    "1. 確認音頻格式為支援的格式 (mp3, wav, m4a 等)\n"
                    "2. 檢查檔案是否完整且未損壞\n"
                    "3. 如果檔案太大 (>25MB)，請先分割檔案\n"
                    "4. 嘗試使用 ffmpeg 轉換為 MP3 格式:\n"
                    "   ffmpeg -i input.m4a -ar 16000 -ac 1 -c:a libmp3lame output.mp3"
                )
            else:
                raise Exception(f"GPT-4o 轉錄失敗: {str(e)}")
    
    def transcription_completed(self, success, result):
        """轉錄完成後的處理"""
        self.transcribe_progress.stop()
        self.transcribe_btn.config(state=tk.NORMAL)
        
        if success:
            self.transcribe_status_var.set(f"轉錄完成: {result}")
            messagebox.showinfo("成功", f"轉錄結果已保存到: {result}")
        else:
            self.transcribe_status_var.set(f"轉錄失敗: {result}")
            messagebox.showerror("錯誤", f"音頻轉錄失敗: {result}")
            
    def setup_slide_ui(self):
        """設置幻燈片捕獲界面"""
        frame = self.slide_frame
        
        # 頂部說明文字
        info_label = tk.Label(
            frame, 
            text="此功能可從視頻文件中捕獲幻燈片，並將其保存為圖片",
            font=("Arial", 12)
        )
        info_label.pack(pady=20)
        
        # 視頻文件選擇區域
        video_frame = tk.Frame(frame)
        video_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(video_frame, text="視頻文件:").pack(side=tk.LEFT, padx=10)
        self.slide_video_entry = tk.Entry(video_frame, width=50)
        self.slide_video_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        
        self.slide_video_browse_btn = tk.Button(
            video_frame, text="瀏覽...", 
            command=lambda: self.browse_file(
                self.slide_video_entry, 
                filetypes=[("視頻文件", "*.mp4 *.avi *.mkv *.mov")]
            )
        )
        self.slide_video_browse_btn.pack(side=tk.LEFT, padx=5)
        
        # 輸出文件夾設置區域
        output_frame = tk.Frame(frame)
        output_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(output_frame, text="輸出文件夾:").pack(side=tk.LEFT, padx=10)
        self.slide_output_entry = tk.Entry(output_frame, width=50)
        self.slide_output_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        
        self.slide_browse_btn = tk.Button(
            output_frame, text="瀏覽...", 
            command=self.browse_slide_folder
        )
        self.slide_browse_btn.pack(side=tk.LEFT, padx=5)
        
        # 相似度閾值設置
        threshold_frame = tk.Frame(frame)
        threshold_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(threshold_frame, text="相似度閾值:").pack(side=tk.LEFT, padx=10)
        
        self.threshold_var = tk.DoubleVar(value=0.95)
        self.threshold_scale = tk.Scale(
            threshold_frame, from_=0.5, to=0.98, 
            variable=self.threshold_var, 
            length=200,
            orient=tk.HORIZONTAL,
            resolution=0.01,
            showvalue=False
        )
        self.threshold_scale.pack(side=tk.LEFT, padx=5)
        
        # 顯示當前閾值
        self.threshold_value_var = tk.StringVar(value="0.95")
        self.threshold_label = tk.Label(
            threshold_frame, 
            textvariable=self.threshold_value_var
        )
        self.threshold_label.pack(side=tk.LEFT, padx=5)
        
        # 更新閾值顯示
        def update_threshold_label(*args):
            value = self.threshold_var.get()
            self.threshold_value_var.set(f"{value:.2f}")
            
        self.threshold_var.trace_add("write", update_threshold_label)
        
        # 說明文字
        tk.Label(
            threshold_frame, 
            text="（值越低檢測越敏感，可能會捕獲較多幻燈片）"
        ).pack(side=tk.LEFT, padx=10)
        
        # 捕獲模式選擇
        mode_frame = tk.Frame(frame)
        mode_frame.pack(fill=tk.X, padx=20, pady=5)
        
        tk.Label(mode_frame, text="捕獲模式:").pack(side=tk.LEFT, padx=5)
        
        self.capture_mode_var = tk.StringVar(value="standard")
        
        tk.Radiobutton(
            mode_frame, text="標準模式", 
            variable=self.capture_mode_var, value="standard"
        ).pack(side=tk.LEFT, padx=5)
        
        tk.Radiobutton(
            mode_frame, text="改進模式（更快更準確）", 
            variable=self.capture_mode_var, value="improved"
        ).pack(side=tk.LEFT, padx=5)
        
        tk.Radiobutton(
            mode_frame, text="超級模式（檢測動畫效果）", 
            variable=self.capture_mode_var, value="ultra"
        ).pack(side=tk.LEFT, padx=5)
        
        tk.Radiobutton(
            mode_frame, text="進階模式（智能分組）", 
            variable=self.capture_mode_var, value="advanced"
        ).pack(side=tk.LEFT, padx=5)
        
        # 狀態顯示
        self.slide_status_var = tk.StringVar(value="準備就緒")
        status_label = tk.Label(
            frame, textvariable=self.slide_status_var,
            font=("Arial", 10), fg="blue"
        )
        status_label.pack(pady=10)
        
        # 進度條
        self.slide_progress = ttk.Progressbar(
            frame, orient="horizontal", length=300, mode="indeterminate"
        )
        self.slide_progress.pack(pady=10)
        
        # 捕獲按鈕
        self.capture_btn = tk.Button(
            frame, text="捕獲幻燈片", command=self.capture_slides,
            bg="#2196F3", fg="white", height=2, width=20
        )
        self.capture_btn.pack(pady=20)
        
    def browse_slide_folder(self):
        """瀏覽並選擇幻燈片輸出文件夾"""
        # 獲取當前輸入框中的路徑
        current_path = self.slide_output_entry.get()
        
        # 如果當前有路徑，從其父目錄開始；否則使用當前工作目錄
        if current_path:
            # 如果路徑存在，使用它；否則使用其父目錄
            if os.path.exists(current_path):
                initial_dir = current_path
            else:
                initial_dir = os.path.dirname(current_path)
                if not os.path.exists(initial_dir):
                    initial_dir = "."
        else:
            initial_dir = "."
            
        folder_path = filedialog.askdirectory(
            initialdir=initial_dir,
            title="選擇保存幻燈片的文件夾"
        )
        if folder_path:
            self.slide_output_entry.delete(0, tk.END)
            self.slide_output_entry.insert(0, folder_path)
    
    def capture_slides(self):
        """捕獲幻燈片按鈕點擊處理"""
        video_path = self.slide_video_entry.get()
        output_folder = self.slide_output_entry.get()
        threshold = self.threshold_var.get()
        
        if not video_path:
            messagebox.showwarning("警告", "請選擇視頻文件")
            return
            
        # 如果未指定輸出文件夾，在視頻所在目錄創建
        if not output_folder:
            video_dir = os.path.dirname(video_path)
            output_folder = os.path.join(video_dir, "slides")
            
            self.slide_output_entry.delete(0, tk.END)
            self.slide_output_entry.insert(0, output_folder)
        
        # 檢查輸出文件夾是否存在，如果不存在則詢問是否創建
        if not os.path.exists(output_folder):
            parent_dir = os.path.dirname(output_folder)
            if parent_dir and not os.path.exists(parent_dir):
                # 父目錄也不存在
                response = messagebox.askyesno(
                    "創建文件夾", 
                    f"文件夾路徑不存在：\n{output_folder}\n\n是否創建此文件夾？"
                )
                if not response:
                    return
            else:
                # 只是最後的文件夾不存在
                response = messagebox.askyesno(
                    "創建文件夾", 
                    f"輸出文件夾不存在：\n{output_folder}\n\n是否創建？"
                )
                if not response:
                    return
        
        # 開始捕獲（在背景線程中運行）
        self.slide_status_var.set("正在捕獲幻燈片...")
        self.slide_progress.start(10)
        self.capture_btn.config(state=tk.DISABLED)
        
        def capture_thread():
            # 根據選擇的模式使用不同的捕獲方法
            capture_mode = self.capture_mode_var.get()
            
            if capture_mode == "improved":
                try:
                    # 嘗試使用改進的捕獲方法
                    from improved_slide_capture import capture_slides_improved
                    success, result = capture_slides_improved(
                        video_path, output_folder, threshold
                    )
                except ImportError:
                    # 如果無法導入改進模組，回退到標準方法
                    self.root.after(0, lambda: self.slide_status_var.set(
                        "改進模組不可用，使用標準模式..."
                    ))
                    success, result = capture_slides_from_video(
                        video_path, output_folder, threshold
                    )
            elif capture_mode == "ultra":
                try:
                    # 嘗試使用超級捕獲方法（檢測動畫）
                    from ultra_slide_capture import capture_slides_ultra
                    self.root.after(0, lambda: self.slide_status_var.set(
                        "使用超級模式，檢測動畫效果..."
                    ))
                    success, result = capture_slides_ultra(
                        video_path, output_folder, threshold
                    )
                except ImportError:
                    # 如果無法導入超級模組，回退到改進方法
                    self.root.after(0, lambda: self.slide_status_var.set(
                        "超級模組不可用，嘗試改進模式..."
                    ))
                    try:
                        from improved_slide_capture import capture_slides_improved
                        success, result = capture_slides_improved(
                            video_path, output_folder, threshold
                        )
                    except ImportError:
                        # 最終回退到標準方法
                        success, result = capture_slides_from_video(
                            video_path, output_folder, threshold
                        )
            elif capture_mode == "advanced":
                try:
                    # 使用進階捕獲方法（智能分組）
                    from slide_capture_advanced import capture_slides_advanced
                    self.root.after(0, lambda: self.slide_status_var.set(
                        "使用進階模式，智能分組..."
                    ))
                    # 進階模式使用兩個閾值：相似度閾值和分組閾值
                    group_threshold = min(threshold + 0.05, 0.95)  # 分組閾值稍高
                    success, result = capture_slides_advanced(
                        video_path, output_folder, threshold, group_threshold
                    )
                except ImportError:
                    # 如果無法導入進階模組，回退到改進方法
                    self.root.after(0, lambda: self.slide_status_var.set(
                        "進階模組不可用，使用改進模式..."
                    ))
                    try:
                        from improved_slide_capture import capture_slides_improved
                        success, result = capture_slides_improved(
                            video_path, output_folder, threshold
                        )
                    except ImportError:
                        # 最終回退到標準方法
                        success, result = capture_slides_from_video(
                            video_path, output_folder, threshold
                        )
            else:
                # 使用標準捕獲方法
                success, result = capture_slides_from_video(
                    video_path, output_folder, threshold
                )
            
            # 在主線程中更新 UI
            self.root.after(0, lambda: self.capture_completed(success, result))
        
        threading.Thread(target=capture_thread).start()
    
    def capture_completed(self, success, result):
        """幻燈片捕獲完成後的處理"""
        self.slide_progress.stop()
        self.capture_btn.config(state=tk.NORMAL)
        
        if success:
            output_folder = result.get("output_folder", "未知文件夾")
            slide_count = result.get("slide_count", 0)
            
            # 構建狀態消息
            status_msg = f"捕獲完成: {slide_count} 張幻燈片"
            
            # 如果有分組信息（進階模式）
            if "group_count" in result:
                group_count = result["group_count"]
                status_msg += f" (分為 {group_count} 組)"
            
            self.slide_status_var.set(status_msg)
            
            # 同時更新處理標籤頁的文件夾路徑
            self.process_folder_entry.delete(0, tk.END)
            self.process_folder_entry.insert(0, output_folder)
            
            # 構建詳細消息
            detail_msg = f"已成功捕獲 {slide_count} 張幻燈片到文件夾: {output_folder}\n"
            
            # 添加進階模式的額外信息
            if "group_count" in result:
                detail_msg += f"幻燈片已分為 {group_count} 組，方便後續處理\n"
                if "metadata_file" in result:
                    detail_msg += f"元數據已保存到: {os.path.basename(result['metadata_file'])}\n"
            
            detail_msg += "\n是否立即處理這些幻燈片？"
            
            # 詢問是否切換到處理標籤頁
            if messagebox.askyesno("成功", detail_msg):
                self.notebook.select(2)  # 切換到處理標籤頁
        else:
            self.slide_status_var.set(f"捕獲失敗: {result}")
            messagebox.showerror("錯誤", f"幻燈片捕獲失敗: {result}")

    def setup_process_ui(self):
        """設置幻燈片處理界面"""
        frame = self.process_frame
        
        # 頂部說明文字
        info_label = tk.Label(
            frame, 
            text="此功能可將捕獲的幻燈片轉換為 PowerPoint 或 Markdown",
            font=("Arial", 12)
        )
        info_label.pack(pady=20)
        
        # 幻燈片文件夾選擇區域
        folder_frame = tk.Frame(frame)
        folder_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(folder_frame, text="幻燈片文件夾:").pack(side=tk.LEFT, padx=10)
        self.process_folder_entry = tk.Entry(folder_frame, width=50)
        self.process_folder_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        
        self.process_browse_btn = tk.Button(
            folder_frame, text="瀏覽...", 
            command=self.browse_process_folder
        )
        self.process_browse_btn.pack(side=tk.LEFT, padx=5)
        
        # 輸出格式選擇
        format_frame = tk.Frame(frame)
        format_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(format_frame, text="輸出格式:").pack(side=tk.LEFT, padx=10)
        
        self.output_format_var = tk.StringVar(value="both")
        
        tk.Radiobutton(
            format_frame, text="僅 PowerPoint", 
            variable=self.output_format_var, value="pptx"
        ).pack(side=tk.LEFT, padx=5)
        
        tk.Radiobutton(
            format_frame, text="僅 Markdown", 
            variable=self.output_format_var, value="markdown"
        ).pack(side=tk.LEFT, padx=5)
        
        tk.Radiobutton(
            format_frame, text="兩者都要", 
            variable=self.output_format_var, value="both"
        ).pack(side=tk.LEFT, padx=5)
        
        # Markdown 處理選項
        md_frame = tk.Frame(frame)
        md_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(md_frame, text="Markdown 處理:").pack(side=tk.LEFT, padx=10)
        
        self.md_process_var = tk.StringVar(value="basic")
        
        tk.Radiobutton(
            md_frame, text="基本處理", 
            variable=self.md_process_var, value="basic"
        ).pack(side=tk.LEFT, padx=5)
        
        tk.Radiobutton(
            md_frame, text="使用 MarkItDown (推薦)", 
            variable=self.md_process_var, value="markitdown"
        ).pack(side=tk.LEFT, padx=5)
        
        # API Key 設置
        api_frame = tk.Frame(frame)
        api_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(api_frame, text="OpenAI API Key:").pack(side=tk.LEFT, padx=10)
        self.api_key_entry = tk.Entry(api_frame, width=50)
        self.api_key_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        
        # 從環境變數獲取 API Key
        api_key_env = os.environ.get("OPENAI_API_KEY", "")
        if api_key_env:
            self.api_key_entry.insert(0, api_key_env)
            
        # 提示文字
        tip_label = tk.Label(
            frame, 
            text="注意: 使用 MarkItDown 可以更好地提取幻燈片中的文字，但需要安裝 markitdown 套件",
            font=("Arial", 10), fg="gray"
        )
        tip_label.pack(pady=5)
        
        # 標題設置
        title_frame = tk.Frame(frame)
        title_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(title_frame, text="文檔標題:").pack(side=tk.LEFT, padx=10)
        self.title_entry = tk.Entry(title_frame, width=50)
        self.title_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        self.title_entry.insert(0, "視頻捕獲的幻燈片")
        
        # 狀態顯示
        self.process_status_var = tk.StringVar(value="準備就緒")
        status_label = tk.Label(
            frame, textvariable=self.process_status_var,
            font=("Arial", 10), fg="blue"
        )
        status_label.pack(pady=10)
        
        # 進度條
        self.process_progress = ttk.Progressbar(
            frame, orient="horizontal", length=300, mode="indeterminate"
        )
        self.process_progress.pack(pady=10)
        
        # 處理按鈕
        self.process_btn = tk.Button(
            frame, text="處理幻燈片", command=self.process_slides,
            bg="#4CAF50", fg="white", height=2, width=20
        )
        self.process_btn.pack(pady=20)
    
    def browse_process_folder(self):
        """瀏覽並選擇幻燈片處理的源文件夾"""
        folder_path = filedialog.askdirectory(
            initialdir=".",
            title="選擇包含幻燈片的文件夾"
        )
        if folder_path:
            self.process_folder_entry.delete(0, tk.END)
            self.process_folder_entry.insert(0, folder_path)
    
    def process_slides(self):
        """處理幻燈片按鈕點擊處理"""
        folder = self.process_folder_entry.get()
        output_format = self.output_format_var.get()
        md_process = self.md_process_var.get()
        api_key = self.api_key_entry.get()
        title = self.title_entry.get()
        
        if not folder:
            messagebox.showwarning("警告", "請選擇幻燈片文件夾")
            return
        
        if not os.path.exists(folder):
            messagebox.showerror("錯誤", f"文件夾不存在: {folder}")
            return
            
        # 確保選擇的是目錄而不是文件
        if not os.path.isdir(folder):
            messagebox.showerror("錯誤", f"選擇的路徑不是目錄: {folder}")
            return
            
        # 檢查文件夾中是否有圖片（排除 macOS 隱藏文件）
        has_images = False
        for filename in os.listdir(folder):
            # 跳過 macOS 隱藏文件
            if filename.startswith('._'):
                continue
            if filename.lower().endswith(('.png', '.jpg', '.jpeg')):
                has_images = True
                break
                
        if not has_images:
            messagebox.showwarning("警告", "選定的文件夾中未找到圖片")
            return
            
        # 開始處理（在背景線程中運行）
        self.process_status_var.set("正在處理幻燈片...")
        self.process_progress.start(10)
        self.process_btn.config(state=tk.DISABLED)
        
        def process_thread():
            # 處理結果
            results = []
            success = True
            
            # 根據選擇處理 PowerPoint
            if output_format in ["pptx", "both"]:
                ppt_success, ppt_result = generate_ppt_from_images(
                    folder, None, title
                )
                success = success and ppt_success
                if ppt_success:
                    results.append(f"PPT: {ppt_result}")
                else:
                    results.append(f"PPT 失敗: {ppt_result}")
            
            # 根據選擇處理 Markdown
            if output_format in ["markdown", "both"]:
                use_markitdown = (md_process == "markitdown")
                md_success, md_result = generate_markdown_from_images(
                    folder, None, title, use_markitdown, api_key
                )
                success = success and md_success
                if md_success:
                    results.append(f"Markdown: {md_result}")
                else:
                    results.append(f"Markdown 失敗: {md_result}")
            
            # 在主線程中更新 UI
            self.root.after(0, lambda: self.processing_completed(success, results))
        
        threading.Thread(target=process_thread).start()
    
    def processing_completed(self, success, results):
        """幻燈片處理完成後的處理"""
        self.process_progress.stop()
        self.process_btn.config(state=tk.NORMAL)
        
        if success:
            self.process_status_var.set("處理完成")
            
            # 顯示處理結果
            result_text = "\n".join(results)
            messagebox.showinfo("成功", f"幻燈片處理完成:\n{result_text}")
        else:
            self.process_status_var.set("處理失敗")
            
            # 顯示處理結果
            result_text = "\n".join(results)
            messagebox.showerror("錯誤", f"幻燈片處理失敗:\n{result_text}")


def main():
    """主函數"""
    # 檢查依賴
    missing_packages = check_dependencies()
    
    if missing_packages:
        print(f"缺少以下依賴: {', '.join(missing_packages)}")
        choice = input("是否自動安裝這些依賴？(y/n): ")
        
        if choice.lower() == 'y':
            if not install_dependencies(missing_packages):
                print("依賴安裝失敗，請手動執行：")
                print(f"pip install {' '.join(missing_packages)}")
                sys.exit(1)
        else:
            print("請手動安裝以下依賴後再運行:")
            print(f"pip install {' '.join(missing_packages)}")
            sys.exit(1)
    
    try:
        # 創建並運行應用
        root = tk.Tk()
        app = VideoAudioProcessor(root)
        
        # 顯示使用提示
        messagebox.showinfo(
            "歡迎使用", 
            "歡迎使用視頻和音頻處理工具\n\n"
            "本工具提供以下功能：\n"
            "1. 從視頻文件中提取音頻\n"
            "2. 從視頻文件中捕獲幻燈片\n"
            "3. 將捕獲的幻燈片轉換為 PowerPoint 或 Markdown\n\n"
            "請選擇相應的標籤頁開始使用"
        )
        
        root.mainloop()
        
    except Exception as e:
        error_msg = traceback.format_exc()
        print(f"啟動失敗: {str(e)}\n{error_msg}")
        messagebox.showerror("啟動錯誤", f"程序啟動失敗:\n{str(e)}")


if __name__ == "__main__":
    main() 