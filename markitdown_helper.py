#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
MarkItDown 輔助工具
提供從圖片轉換到 Markdown 的功能
"""

import os
import base64
import tempfile
import traceback
from typing import List, Tuple, Dict, Any, Optional


def _convert_heic_to_jpeg(heic_path: str) -> str:
    """
    將 HEIC/HEIF 圖片轉換為 JPEG 格式的臨時文件

    參數:
        heic_path: HEIC 圖片路徑

    返回:
        轉換後的 JPEG 臨時文件路徑
    """
    try:
        import pillow_heif
        pillow_heif.register_heif_opener()
    except ImportError:
        raise ImportError(
            "需要安裝 pillow-heif 來處理 HEIC 格式圖片。"
            "請執行: pip install pillow-heif"
        )

    from PIL import Image
    img = Image.open(heic_path)
    if img.mode != 'RGB':
        img = img.convert('RGB')

    temp_file = tempfile.NamedTemporaryFile(suffix='.jpg', delete=False)
    img.save(temp_file.name, 'JPEG', quality=95)
    temp_file.close()
    return temp_file.name


def convert_images_to_markdown(
    image_paths: List[str],
    output_file: str,
    title: str = "圖片內容分析",
    use_llm: bool = False,
    api_key: Optional[str] = None,
    model: str = "gpt-4o-mini"
) -> Tuple[bool, str, Dict[str, Any]]:
    """
    將圖片文件轉換為 Markdown 文件
    
    參數:
        image_paths: 圖片文件路徑列表
        output_file: 輸出的 Markdown 文件路徑
        title: Markdown 文件標題
        use_llm: 是否使用 LLM 進行圖片文字識別與分析
        api_key: OpenAI API Key，只有當 use_llm=True 時才有效
        model: 使用的 LLM 模型，只有當 use_llm=True 時才有效
        
    返回:
        success: 是否成功
        output_file: 輸出文件路徑
        info: 包含轉換統計信息的字典
    """
    heic_temp_files = []  # 追蹤需要清理的臨時文件
    try:
        # 過濾掉 macOS 的隱藏文件和非圖片文件
        valid_image_paths = []
        skipped_files = []
        supported_formats = {'.png', '.jpg', '.jpeg', '.gif', '.webp', '.heic', '.heif'}
        heic_formats = {'.heic', '.heif'}

        for img_path in image_paths:
            basename = os.path.basename(img_path)
            # 跳過以 ._ 開頭的 macOS 隱藏文件
            if basename.startswith('._'):
                skipped_files.append(img_path)
                print(f"跳過 macOS 隱藏文件: {basename}")
                continue

            # 檢查文件擴展名
            ext = os.path.splitext(basename)[1].lower()
            if ext not in supported_formats:
                skipped_files.append(img_path)
                print(f"跳過不支持的文件格式: {basename}")
                continue

            # HEIC/HEIF 格式需要轉換為 JPEG
            if ext in heic_formats:
                try:
                    print(f"轉換 HEIC 圖片: {basename}")
                    converted_path = _convert_heic_to_jpeg(img_path)
                    heic_temp_files.append(converted_path)
                    valid_image_paths.append(converted_path)
                except ImportError as e:
                    print(str(e))
                    skipped_files.append(img_path)
                    continue
                except Exception as e:
                    print(f"轉換 HEIC 圖片失敗 {basename}: {e}")
                    skipped_files.append(img_path)
                    continue
            else:
                valid_image_paths.append(img_path)

        if not valid_image_paths:
            return False, "", {
                "error": "沒有有效的圖片文件可處理",
                "skipped_files": skipped_files,
                "success": False
            }
        
        info = {
            "success": True, 
            "processed_images": len(valid_image_paths),
            "skipped_files": len(skipped_files),
            "total_files": len(image_paths)
        }
        
        # 如果需要使用 LLM 進行圖片識別和分析
        if use_llm and api_key:
            try:
                import openai
                
                # 設置 API Key
                openai.api_key = api_key
                
                print(f"使用 {model} 模型分析 {len(valid_image_paths)} 張圖片...")
                with open(output_file, 'w', encoding='utf-8') as f:
                    f.write(f"# {title}\n\n")
                    
                    for i, img_path in enumerate(valid_image_paths):
                        # 添加標題和圖片
                        slide_num = i + 1
                        f.write(f"## 幻燈片 {slide_num}\n\n")
                        
                        # 獲取相對路徑
                        try:
                            rel_path = os.path.relpath(
                                img_path, 
                                os.path.dirname(output_file)
                            )
                        except Exception:
                            rel_path = img_path
                            
                        f.write(f"![幻燈片 {slide_num}]({rel_path})\n\n")
                        
                        # 使用 OpenAI 視覺模型分析圖片
                        print(
                            f"分析圖片 {slide_num}/{len(valid_image_paths)}: "
                            f"{os.path.basename(img_path)}"
                        )
                        
                        try:
                            with open(img_path, "rb") as img_file:
                                # 創建 OpenAI 客戶端
                                client = openai.OpenAI(api_key=api_key)
                                
                                # 使用 Vision API 分析圖片內容
                                content_system = (
                                    "你是一個幻燈片分析專家。請識別並提取圖片中所有可見的"
                                    "文本內容，同時分析圖片中的圖表、表格和其他視覺元素。"
                                    "以結構化的Markdown格式返回內容，保持原始格式和層次結構。"
                                )
                                
                                encoded_img = base64.b64encode(
                                    img_file.read()
                                ).decode('utf-8')
                                
                                response = client.chat.completions.create(
                                    model=model,
                                    messages=[
                                        {
                                            "role": "system",
                                            "content": content_system
                                        },
                                        {
                                            "role": "user",
                                            "content": [
                                                {
                                                    "type": "text", 
                                                    "text": "請分析這張幻燈片圖片並提取其中的內容："
                                                },
                                                {
                                                    "type": "image_url",
                                                    "image_url": {
                                                        "url": f"data:image/jpeg;base64,{encoded_img}"
                                                    }
                                                }
                                            ]
                                        }
                                    ],
                                    max_tokens=1000
                                )
                                
                                # 寫入分析結果
                                f.write(
                                    f"{response.choices[0].message.content}\n\n"
                                )
                                f.write("---\n\n")
                        except Exception as e:
                            error_msg = str(e)
                            print(f"分析圖片 {img_path} 時出錯: {error_msg}")
                            f.write(f"*分析圖片時出錯: {error_msg}*\n\n")
                            f.write("---\n\n")
                
                info["llm_used"] = True
                info["model"] = model
                return True, output_file, info
                
            except ImportError:
                print("OpenAI 模組未安裝，將使用基本方法")
                pass
            except Exception as e:
                print(f"使用 LLM 分析圖片時出錯: {str(e)}")
                traceback.print_exc()
                info["error"] = str(e)
                info["llm_used"] = False
                pass
        
        # 使用基本方法生成 Markdown
        print("使用基本方法生成 Markdown...")
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(f"# {title}\n\n")
            
            for i, img_path in enumerate(valid_image_paths):
                # 獲取相對路徑
                try:
                    rel_path = os.path.relpath(
                        img_path, 
                        os.path.dirname(output_file)
                    )
                except Exception:
                    rel_path = img_path
                    
                # 添加幻燈片標題和圖片
                slide_num = i + 1
                f.write(f"## 幻燈片 {slide_num}\n\n")
                f.write(f"![幻燈片 {slide_num}]({rel_path})\n\n")
                f.write("---\n\n")
        
        info["llm_used"] = False
        return True, output_file, info
        
    except Exception as e:
        error_msg = str(e)
        print(f"生成 Markdown 時出錯: {error_msg}")
        traceback.print_exc()
        return False, "", {"error": error_msg, "success": False}
    finally:
        # 清理 HEIC 轉換的臨時文件
        for temp_path in heic_temp_files:
            try:
                if os.path.exists(temp_path):
                    os.unlink(temp_path)
            except OSError:
                pass


def process_images_to_ppt(
    image_dir: str,
    output_ppt: str,
    title: str = "幻燈片簡報",
    use_llm: bool = False,
    api_key: Optional[str] = None,
    model: str = "gpt-4o-mini"
) -> bool:
    """
    將文件夾中的圖片轉換為 PowerPoint 文件
    
    參數:
        image_dir: 包含幻燈片圖片的文件夾
        output_ppt: 輸出的 PowerPoint 文件路徑
        title: 演示文稿標題
        use_llm: 是否使用 LLM 添加圖片描述（尚未實現）
        api_key: OpenAI API Key，只有當 use_llm=True 時才有效（尚未實現）
        model: 使用的 LLM 模型，只有當 use_llm=True 時才有效（尚未實現）
        
    返回:
        success: 是否成功
    """
    heic_temp_files_ppt = []
    try:
        from pptx import Presentation
        from pptx.util import Inches
        from slide_sort import sorted_image_paths

        # 先按「檔名時間標記 → 檔案建立時間」排序，再對 HEIC 做轉換（保留順序）
        source_paths = sorted_image_paths(
            image_dir,
            extensions=('.png', '.jpg', '.jpeg', '.heic', '.heif'),
        )

        image_files = []
        for src in source_paths:
            ext = os.path.splitext(src)[1].lower()
            if ext in ('.heic', '.heif'):
                try:
                    print(f"轉換 HEIC 圖片: {os.path.basename(src)}")
                    converted = _convert_heic_to_jpeg(src)
                    heic_temp_files_ppt.append(converted)
                    image_files.append(converted)
                except Exception as e:
                    print(f"轉換 HEIC 圖片失敗 {os.path.basename(src)}: {e}")
            else:
                image_files.append(src)

        if not image_files:
            print("未找到圖片文件")
            return False

        # 創建新的 PowerPoint 演示文稿
        prs = Presentation()
        
        # 設置幻燈片尺寸為 16:9
        prs.slide_width = Inches(10)  # 16:9 寬度
        prs.slide_height = Inches(5.625)  # 16:9 高度
        
        # 獲取幻燈片尺寸
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # 為每張圖片創建一個幻燈片（不添加標題幻燈片）
        blank_slide_layout = prs.slide_layouts[6]  # 空白布局
        
        for img_path in image_files:
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 添加圖片
            left = top = Inches(0)
            _ = slide.shapes.add_picture(
                img_path, left, top, width=slide_width, height=slide_height
            )
            
        # 保存演示文稿
        prs.save(output_ppt)
        print(f"已成功生成 PowerPoint: {output_ppt}")

        return True

    except Exception as e:
        error_msg = str(e)
        print(f"生成 PowerPoint 時出錯: {error_msg}")
        traceback.print_exc()
        return False
    finally:
        # 清理 HEIC 轉換的臨時文件
        for temp_path in heic_temp_files_ppt:
            try:
                if os.path.exists(temp_path):
                    os.unlink(temp_path)
            except OSError:
                pass