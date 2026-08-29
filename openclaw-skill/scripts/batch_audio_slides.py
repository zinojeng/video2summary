#!/usr/bin/env python3
from __future__ import annotations
"""
batch_audio_slides.py — 批次處理 ~/audio/ 所有子資料夾

只做兩件事（不合併）：
  1. 音訊轉錄 → .txt
  2. 投影片分析 → slides_analysis.md

用法：
  python batch_audio_slides.py                    # 處理所有資料夾
  python batch_audio_slides.py --folder "glp1"    # 只處理指定資料夾
  python batch_audio_slides.py --dry-run           # 預覽不執行
  python batch_audio_slides.py --language en       # 指定轉錄語言
  python batch_audio_slides.py --retranscribe      # 強制重新轉錄（忽略已有 .txt）
"""

import argparse
import json
import os
import subprocess
import sys
from pathlib import Path
from datetime import datetime

# 讓 openclaw-skill/scripts 也能 import 專案根目錄的共用設定
sys.path.insert(0, str(Path(__file__).resolve().parents[2]))
from model_config import OPENAI_VISION

AUDIO_DIR = Path.home() / "audio"
SKILL_DIR = Path.home() / ".openclaw" / "workspace" / "skills" / "audio-transcribe" / "scripts"
PYTHON = SKILL_DIR / "venv" / "bin" / "python3"
TRANSCRIBER = SKILL_DIR / "gpt4o_transcribe_improved.py"

AUDIO_EXTENSIONS = {
    "mp3", "m4a", "wav", "ogg", "opus", "flac", "aac", "wma",
    "mp4", "mov", "mkv", "avi", "webm"
}
IMAGE_EXTENSIONS = {"png", "jpg", "jpeg", "gif", "webp", "heic", "heif"}
SLIDE_EXTENSIONS = IMAGE_EXTENSIONS | {"pptx", "pdf"}

SKIP_DIRS = {"done", "test-slides", ".Trash"}


def log(msg):
    ts = datetime.now().strftime("%H:%M:%S")
    print(f"[{ts}] {msg}", flush=True)


def ext(p: Path) -> str:
    return p.suffix.lstrip(".").lower()


def find_files(folder: Path, extensions: set) -> list[Path]:
    return sorted([
        f for f in folder.iterdir()
        if f.is_file() and ext(f) in extensions and not f.name.startswith(".")
    ])


def get_env():
    env = os.environ.copy()
    if "/opt/homebrew/bin" not in env.get("PATH", ""):
        env["PATH"] = f"/opt/homebrew/bin:{env.get('PATH', '/usr/bin:/bin')}"
    env_file = SKILL_DIR / ".env"
    if env_file.exists():
        for line in env_file.read_text().splitlines():
            line = line.strip()
            if line.startswith("export "):
                line = line[7:]
            if "=" in line and not line.startswith("#"):
                key, _, val = line.partition("=")
                val = val.strip().strip('"')
                env.setdefault(key.strip(), val)
    return env


def has_transcript(folder: Path):
    """檢查資料夾是否已有轉錄檔"""
    for f in folder.iterdir():
        if f.is_file() and f.suffix == ".txt" and not f.name.startswith("."):
            # 排除 slides_analysis 等非轉錄檔
            if "analysis" not in f.name.lower() and "note" not in f.name.lower():
                return f
    # 也檢查 ~/audio/ 下同名 .txt
    parent_txt = folder.parent / f"{folder.name}.txt"
    if parent_txt.exists():
        return parent_txt
    return None


def has_slides_analysis(folder: Path) -> bool:
    return (folder / "slides_analysis.md").exists()


def transcribe_audio(audio_path: Path, output_dir: Path, language: str = None) -> bool:
    """轉錄音訊檔"""
    stem = audio_path.stem
    output_file = output_dir / f"{stem}.txt"

    cmd = [
        str(PYTHON), str(TRANSCRIBER),
        str(audio_path),
        "--format", "text",
        "--output", str(output_file),
        "--model", "gpt-transcribe",
        "--max-segment-seconds", "300",  # 5 分鐘分段，降低 repetition loop
    ]
    if language:
        cmd += ["--language", language]

    log(f"  🎙️ 轉錄中: {audio_path.name}")
    log(f"     → {output_file.name}")

    result = subprocess.run(cmd, env=get_env(), timeout=1800)  # 30 min timeout
    if result.returncode == 0 and output_file.exists():
        size = output_file.stat().st_size
        log(f"  ✅ 轉錄完成 ({size / 1024:.1f} KB)")
        return True
    else:
        log(f"  ❌ 轉錄失敗 (exit code: {result.returncode})")
        return False


def analyze_slides(folder: Path) -> bool:
    """分析投影片（使用 markitdown_helper 或 OpenAI Vision）"""
    output_file = folder / "slides_analysis.md"

    # 收集圖片
    images = find_files(folder, IMAGE_EXTENSIONS)
    pptx_files = [f for f in folder.iterdir() if f.is_file() and ext(f) == "pptx"]
    pdf_files = [f for f in folder.iterdir() if f.is_file() and ext(f) == "pdf"]

    total = len(images) + len(pptx_files) + len(pdf_files)
    if total == 0:
        return False

    log(f"  🖼️ 分析投影片: {len(images)} 圖片, {len(pptx_files)} PPTX, {len(pdf_files)} PDF")

    env = get_env()
    api_key = env.get("OPENAI_API_KEY", "")

    if not api_key:
        log(f"  ❌ 缺少 OPENAI_API_KEY")
        return False

    # 嘗試用 markitdown_helper
    try:
        if str(SKILL_DIR) not in sys.path:
            sys.path.insert(0, str(SKILL_DIR))
        from markitdown_helper import convert_images_to_markdown

        # 先處理 PPTX/PDF → 提取圖片
        all_images = list(images)
        temp_dir = folder / "_temp_extracted"

        for pptx in pptx_files:
            temp_dir.mkdir(exist_ok=True)
            extracted = extract_images_from_pptx(pptx, temp_dir)
            all_images.extend(extracted)

        for pdf in pdf_files:
            temp_dir.mkdir(exist_ok=True)
            extracted = extract_images_from_pdf(pdf, temp_dir)
            all_images.extend(extracted)

        if not all_images:
            log(f"  ⚠️ 無可分析的圖片")
            return False

        log(f"  🔍 共 {len(all_images)} 張圖片待分析")

        success, out_path, info = convert_images_to_markdown(
            image_paths=[str(p) for p in all_images],
            output_file=str(output_file),
            title=f"{folder.name} 投影片分析",
            use_llm=True,
            api_key=api_key,
            model=OPENAI_VISION
        )

        # 清理暫存
        if temp_dir.exists():
            import shutil
            shutil.rmtree(temp_dir, ignore_errors=True)

        if success:
            log(f"  ✅ 投影片分析完成 ({info.get('processed_images', 0)} 張)")
            return True
        else:
            log(f"  ❌ 分析失敗: {info.get('error', 'unknown')}")
            return False

    except ImportError:
        log(f"  ⚠️ markitdown_helper 未找到，使用 OpenAI Vision 直接分析")
        return analyze_slides_direct(folder, all_images if 'all_images' in dir() else images, output_file, api_key)


def extract_images_from_pptx(pptx_path: Path, output_dir: Path) -> list[Path]:
    """從 PPTX 提取圖片"""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(str(pptx_path))
        images = []
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    ext_name = image.content_type.split("/")[-1]
                    if ext_name == "jpeg":
                        ext_name = "jpg"
                    img_path = output_dir / f"pptx_slide{i+1:03d}.{ext_name}"
                    img_path.write_bytes(image.blob)
                    images.append(img_path)
        log(f"    📦 從 {pptx_path.name} 提取 {len(images)} 張圖片")
        return images
    except Exception as e:
        log(f"    ⚠️ PPTX 提取失敗: {e}")
        return []


def extract_images_from_pdf(pdf_path: Path, output_dir: Path) -> list[Path]:
    """從 PDF 提取圖片"""
    try:
        from pdf2image import convert_from_path
        pages = convert_from_path(str(pdf_path), dpi=150)
        images = []
        for i, page in enumerate(pages):
            img_path = output_dir / f"pdf_page{i+1:03d}.jpg"
            page.save(str(img_path), "JPEG")
            images.append(img_path)
        log(f"    📦 從 {pdf_path.name} 提取 {len(images)} 頁")
        return images
    except Exception as e:
        log(f"    ⚠️ PDF 提取失敗: {e}")
        return []


def analyze_slides_direct(folder, images, output_file, api_key):
    """直接用 OpenAI Vision API 分析"""
    try:
        import base64
        from openai import OpenAI

        client = OpenAI(api_key=api_key)

        with open(output_file, "w", encoding="utf-8") as f:
            f.write(f"# {folder.name} 投影片分析\n\n")
            for i, img_path in enumerate(images):
                log(f"    分析 {i+1}/{len(images)}: {img_path.name}")
                f.write(f"## 投影片 {i+1}\n\n")

                # HEIC 轉換
                actual_path = img_path
                if ext(img_path) in ("heic", "heif"):
                    try:
                        from pillow_heif import register_heif_opener
                        register_heif_opener()
                        from PIL import Image
                        jpg_path = img_path.with_suffix(".jpg")
                        Image.open(img_path).convert("RGB").save(jpg_path, "JPEG", quality=85)
                        actual_path = jpg_path
                    except Exception:
                        pass

                with open(actual_path, "rb") as img_file:
                    encoded = base64.b64encode(img_file.read()).decode("utf-8")

                mime = "image/jpeg" if ext(actual_path) in ("jpg", "jpeg") else f"image/{ext(actual_path)}"
                response = client.chat.completions.create(
                    model=OPENAI_VISION,
                    messages=[
                        {"role": "system", "content":
                            "你是一個醫學投影片分析專家。請識別並提取圖片中所有可見的"
                            "文本內容，同時分析圖表、表格和視覺元素。"
                            "以結構化的 Markdown 格式返回，使用繁體中文。"},
                        {"role": "user", "content": [
                            {"type": "text", "text": "請分析這張投影片並提取內容："},
                            {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{encoded}"}}
                        ]}
                    ],
                    max_completion_tokens=1500
                )
                f.write(f"{response.choices[0].message.content}\n\n---\n\n")

        log(f"  ✅ 投影片分析完成: {output_file.name}")
        return True
    except Exception as e:
        log(f"  ❌ 分析失敗: {e}")
        return False


def scan_and_process(target_folder=None, dry_run=False, language=None, retranscribe=False):
    """掃描並處理所有資料夾"""
    if not AUDIO_DIR.exists():
        log(f"❌ {AUDIO_DIR} 不存在")
        return

    subfolders = sorted([
        d for d in AUDIO_DIR.iterdir()
        if d.is_dir() and not d.name.startswith(".") and d.name not in SKIP_DIRS
    ])

    if target_folder:
        subfolders = [d for d in subfolders if d.name == target_folder]
        if not subfolders:
            log(f"❌ 找不到資料夾: {target_folder}")
            return

    log(f"📂 掃描 {AUDIO_DIR}，共 {len(subfolders)} 個資料夾")
    print()

    stats = {"transcribed": 0, "analyzed": 0, "skipped": 0, "failed": 0}

    for folder in subfolders:
        audio_files = find_files(folder, AUDIO_EXTENSIONS)
        slide_files = find_files(folder, SLIDE_EXTENSIONS)
        existing_txt = has_transcript(folder)
        existing_analysis = has_slides_analysis(folder)

        # 判斷需要做什麼
        need_transcribe = bool(audio_files) and (not existing_txt or retranscribe)
        need_analyze = bool(slide_files) and not existing_analysis

        if not need_transcribe and not need_analyze:
            status = []
            if existing_txt:
                status.append("✅ 已轉錄")
            if existing_analysis:
                status.append("✅ 已分析")
            if not audio_files and not slide_files:
                status.append("(空)")
            log(f"⏭️  {folder.name}/ — {', '.join(status)}")
            stats["skipped"] += 1
            continue

        log(f"📂 {folder.name}/")
        if audio_files:
            log(f"   🎵 音訊: {len(audio_files)} 個 {'(需轉錄)' if need_transcribe else '(已完成)'}")
        if slide_files:
            log(f"   🖼️ 投影片: {len(slide_files)} 個 {'(需分析)' if need_analyze else '(已完成)'}")

        if dry_run:
            continue

        # 1. 轉錄
        if need_transcribe:
            try:
                audio_file = audio_files[0]
                if transcribe_audio(audio_file, folder, language):
                    stats["transcribed"] += 1
                else:
                    stats["failed"] += 1
            except Exception as e:
                log(f"  ❌ 轉錄失敗: {e}")
                stats["failed"] += 1

        # 2. 投影片分析
        if need_analyze:
            try:
                if analyze_slides(folder):
                    stats["analyzed"] += 1
                else:
                    stats["failed"] += 1
            except Exception as e:
                log(f"  ❌ 分析失敗: {e}")
                stats["failed"] += 1

        print()

    # 統計
    print("\n" + "=" * 50)
    log(f"📊 處理完成")
    log(f"   🎙️ 轉錄: {stats['transcribed']} 個")
    log(f"   🖼️ 分析: {stats['analyzed']} 個")
    log(f"   ⏭️ 跳過: {stats['skipped']} 個")
    if stats["failed"]:
        log(f"   ❌ 失敗: {stats['failed']} 個")


def main():
    parser = argparse.ArgumentParser(description="批次處理音訊轉錄 + 投影片分析（不合併）")
    parser.add_argument("--folder", type=str, help="只處理指定資料夾")
    parser.add_argument("--dry-run", action="store_true", help="預覽不執行")
    parser.add_argument("--language", type=str, help="轉錄語言 (例如: en, zh)")
    parser.add_argument("--retranscribe", action="store_true", help="強制重新轉錄")
    args = parser.parse_args()

    scan_and_process(
        target_folder=args.folder,
        dry_run=args.dry_run,
        language=args.language,
        retranscribe=args.retranscribe,
    )


if __name__ == "__main__":
    main()
