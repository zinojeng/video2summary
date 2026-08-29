#!/usr/bin/env python3
from __future__ import annotations

"""
folder_processor.py — ~/audio/ 資料夾監控處理器

監控 ~/audio/ 資料夾，自動處理音訊與投影片：

規則：
  - 檔案（音訊）→ 轉錄 → 整理筆記
  - 子資料夾（音訊+投影片）→ 轉錄 + 投影片分析 → 合併筆記

資料夾結構：
  ~/audio/
  ├── Pit-1.m4a                ← 單一音訊檔 → 僅轉錄+筆記
  ├── HAE_黃教授/              ← 子資料夾 = 音訊+投影片合併
  │   ├── HAE_黃教授.m4a       ← 音訊檔
  │   ├── slide_01.jpg         ← 投影片（圖片/PPTX/PDF 皆可）
  │   ├── slides.pptx
  │   └── slides.pdf
  └── ADA_2026/
      ├── recording.m4a
      └── slides.pdf

投影片可以晚到：先產出 transcript-only 筆記，投影片到了自動重跑 merged 版。
處理完的單一音訊檔移至 ~/audio/done/。

用法：
  python folder_processor.py                 # 掃描一次
  python folder_processor.py --daemon        # 持續監控
  python folder_processor.py --interval 30   # 30秒一輪
"""

import argparse
import json
import os
import subprocess
import sys
import time
import shutil
from pathlib import Path
from datetime import datetime

# 讓 openclaw-skill/scripts 也能 import 專案根目錄的共用設定
sys.path.insert(0, str(Path(__file__).resolve().parents[2]))
from model_config import OPENAI_VISION

# ===== 設定 =====
AUDIO_DIR = Path.home() / "audio"
DONE_DIR = AUDIO_DIR / "done"
STATE_FILE = AUDIO_DIR / ".processor_state.json"
POLL_INTERVAL = 60  # seconds

# 支援的檔案格式
AUDIO_EXTENSIONS = {
    "mp3", "m4a", "wav", "ogg", "opus", "flac", "aac", "wma",
    "mp4", "mov", "mkv", "avi", "webm"
}
IMAGE_EXTENSIONS = {"png", "jpg", "jpeg", "gif", "webp", "heic", "heif"}
SLIDE_EXTENSIONS = IMAGE_EXTENSIONS | {"pptx", "pdf"}

# Skill 路徑
SKILL_DIR = Path.home() / ".openclaw" / "workspace" / "skills" / "audio-transcribe" / "scripts"
PYTHON = SKILL_DIR / "venv" / "bin" / "python3"
TRANSCRIBER = SKILL_DIR / "gpt4o_transcribe_clawbot.py"
SUMMARIZER = SKILL_DIR / "summarize_transcript.py"

# Telegram
TG_BOT_TOKEN = os.environ.get("TG_BOT_TOKEN", "")
TG_CHAT_ID = os.environ.get("TG_CHAT_ID", "")


def log(msg: str):
    ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    print(f"[{ts}] {msg}", flush=True)


def send_telegram(message: str):
    """發送 Telegram 通知（使用 urllib，避免 curl 編碼問題）"""
    if not TG_BOT_TOKEN or not TG_CHAT_ID:
        return
    import urllib.request
    import urllib.parse
    try:
        url = f"https://api.telegram.org/bot{TG_BOT_TOKEN}/sendMessage"
        data = urllib.parse.urlencode({
            "chat_id": TG_CHAT_ID,
            "text": message,
            "parse_mode": "Markdown",
        }).encode("utf-8")
        req = urllib.request.Request(url, data=data, method="POST")
        with urllib.request.urlopen(req, timeout=15) as resp:
            body = resp.read().decode("utf-8")
            if '"ok":false' in body:
                log(f"  ⚠️ Telegram API 回傳錯誤: {body[:200]}")
    except Exception as e:
        log(f"  ⚠️ Telegram 發送失敗: {e}")


def setup_dirs():
    """建立工作目錄"""
    AUDIO_DIR.mkdir(parents=True, exist_ok=True)
    DONE_DIR.mkdir(parents=True, exist_ok=True)


def load_state() -> dict:
    """載入處理狀態"""
    if STATE_FILE.exists():
        try:
            return json.loads(STATE_FILE.read_text())
        except (json.JSONDecodeError, OSError):
            pass
    return {"processed_folders": {}, "processed_audio": []}


def save_state(state: dict):
    """儲存處理狀態"""
    STATE_FILE.write_text(json.dumps(state, ensure_ascii=False, indent=2))


def get_env() -> dict:
    """取得帶有 API key 和 PATH 的環境變數"""
    env = os.environ.copy()
    if "/opt/homebrew/bin" not in env.get("PATH", ""):
        env["PATH"] = f"/opt/homebrew/bin:{env.get('PATH', '/usr/bin:/bin')}"
    if "OPENAI_API_KEY" not in env:
        env_file = SKILL_DIR / ".env"
        if env_file.exists():
            for line in env_file.read_text().splitlines():
                line = line.strip()
                if line.startswith("export "):
                    line = line[7:]
                if "=" in line and not line.startswith("#"):
                    key, _, val = line.partition("=")
                    val = val.strip().strip('"')
                    env.setdefault(key.strip(), val)
    return env


def _ext(p: Path) -> str:
    return p.suffix.lstrip(".").lower()


def find_audio_files(folder: Path) -> list[Path]:
    """找出資料夾中的音訊檔"""
    return sorted([
        f for f in folder.iterdir()
        if f.is_file() and _ext(f) in AUDIO_EXTENSIONS and not f.name.startswith(".")
    ])


def find_slide_files(folder: Path) -> list[Path]:
    """找出資料夾中的投影片檔（圖片/PPTX/PDF）"""
    return sorted([
        f for f in folder.iterdir()
        if f.is_file() and _ext(f) in SLIDE_EXTENSIONS and not f.name.startswith(".")
    ])


def find_image_files(folder: Path) -> list[Path]:
    """找出資料夾中的圖片檔"""
    return sorted([
        f for f in folder.iterdir()
        if f.is_file() and _ext(f) in IMAGE_EXTENSIONS and not f.name.startswith(".")
    ])


# ===== A. 音訊處理 =====

def transcribe_audio(audio_path: Path, output_dir: Path) -> Path | None:
    """轉錄音訊檔案"""
    stem = audio_path.stem
    output_file = output_dir / f"{stem}.txt"

    log(f"  🎙️ 開始轉錄: {audio_path.name}")

    env = get_env()
    result = subprocess.run(
        [str(PYTHON), str(TRANSCRIBER),
         str(audio_path),
         "--language", "zh",
         "--format", "text",
         "--output", str(output_file)],
        capture_output=True, text=True, timeout=3600,
        stdin=subprocess.DEVNULL,
        env=env, cwd=str(SKILL_DIR)
    )

    if result.returncode == 0 and output_file.exists():
        chars = len(output_file.read_text())
        log(f"  ✅ 轉錄完成: {output_file.name} ({chars} 字)")
        return output_file
    else:
        log(f"  ❌ 轉錄失敗: {result.stderr[-500:] if result.stderr else 'unknown'}")
        return None


def summarize_transcript(transcript_path: Path, slides_md: Path | None = None) -> Path | None:
    """產生整理筆記，可選合併 slides 分析"""
    stem = transcript_path.stem
    output_dir = transcript_path.parent

    if slides_md and slides_md.exists():
        output_file = output_dir / f"{stem}_merged_notes.md"
        log(f"  📝 開始合併筆記: {stem} + slides")
        extra_args = ["--slides", str(slides_md), "--style", "merged"]
    else:
        output_file = output_dir / f"{stem}_detailed_notes.md"
        log(f"  📝 開始整理筆記: {stem}")
        extra_args = ["--style", "medical"]

    env = get_env()
    result = subprocess.run(
        [str(PYTHON), str(SUMMARIZER),
         str(transcript_path),
         "--provider", "openai",
         "--output", str(output_file)] + extra_args,
        capture_output=True, text=True, timeout=600,
        stdin=subprocess.DEVNULL,
        env=env, cwd=str(SKILL_DIR)
    )

    if result.returncode == 0 and output_file.exists():
        chars = len(output_file.read_text())
        log(f"  ✅ 筆記完成: {output_file.name} ({chars} 字)")
        return output_file
    else:
        log(f"  ❌ 筆記整理失敗: {result.stderr[-500:] if result.stderr else 'unknown'}")
        return None


def convert_md_to_docx(md_path: Path) -> Path | None:
    """將 .md 轉成 .docx"""
    docx_path = md_path.with_suffix(".docx")

    # 方法1：docx-js
    converter_script = SKILL_DIR / "md_to_docx.js"
    if converter_script.exists():
        try:
            result = subprocess.run(
                ["/opt/homebrew/bin/node", str(converter_script), str(md_path), str(docx_path)],
                capture_output=True, text=True, timeout=120,
                stdin=subprocess.DEVNULL,
                env={**os.environ, "NODE_PATH": "/opt/homebrew/lib/node_modules"}
            )
            if result.returncode == 0 and docx_path.exists():
                log(f"  ✅ DOCX 轉換完成: {docx_path.name}")
                return docx_path
        except FileNotFoundError:
            pass

    # 方法2：pandoc
    try:
        result = subprocess.run(
            ["pandoc", str(md_path), "-o", str(docx_path),
             "--from", "markdown", "--to", "docx"],
            capture_output=True, text=True, timeout=60,
            stdin=subprocess.DEVNULL
        )
        if result.returncode == 0 and docx_path.exists():
            log(f"  ✅ DOCX 轉換完成 (pandoc): {docx_path.name}")
            return docx_path
    except FileNotFoundError:
        pass

    log(f"  ⚠️ DOCX 轉換失敗")
    return None


# ===== B. 投影片分析 =====

def extract_images_from_pptx(pptx_path: Path, output_dir: Path) -> list[Path]:
    """從 PPTX 中提取所有投影片圖片"""
    images = []
    try:
        from pptx import Presentation
        from PIL import Image
        import io

        prs = Presentation(str(pptx_path))
        for i, slide in enumerate(prs.slides):
            # 嘗試從 slide 的 shapes 中找圖片
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    img_data = shape.image.blob
                    img = Image.open(io.BytesIO(img_data))
                    img_path = output_dir / f"pptx_slide_{i+1:03d}.jpg"
                    img.convert("RGB").save(str(img_path), "JPEG", quality=95)
                    images.append(img_path)
                    break  # 一個 slide 取一張主要圖片

        # 如果沒從 shapes 找到，用整頁截圖方式（需要 libreoffice）
        if not images:
            log(f"  ℹ️ PPTX 無嵌入圖片，嘗試用 LibreOffice 轉圖片...")
            result = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf",
                 "--outdir", str(output_dir), str(pptx_path)],
                capture_output=True, text=True, timeout=120,
                stdin=subprocess.DEVNULL
            )
            pdf_path = output_dir / pptx_path.with_suffix(".pdf").name
            if pdf_path.exists():
                images = extract_images_from_pdf(pdf_path, output_dir)
                pdf_path.unlink()  # 清理臨時 PDF

        log(f"  📊 從 PPTX 提取 {len(images)} 張投影片")
    except Exception as e:
        log(f"  ❌ PPTX 處理失敗: {e}")
    return images


def extract_images_from_pdf(pdf_path: Path, output_dir: Path) -> list[Path]:
    """從 PDF 中提取頁面為圖片"""
    images = []
    try:
        # 方法1: pdf2image (poppler)
        from pdf2image import convert_from_path
        pages = convert_from_path(str(pdf_path), dpi=200)
        for i, page in enumerate(pages):
            img_path = output_dir / f"pdf_page_{i+1:03d}.jpg"
            page.save(str(img_path), "JPEG", quality=95)
            images.append(img_path)
        log(f"  📄 從 PDF 提取 {len(images)} 頁")
    except ImportError:
        # 方法2: 用 sips (macOS) 或 convert
        log(f"  ⚠️ pdf2image 未安裝，嘗試替代方案...")
        try:
            result = subprocess.run(
                ["sips", "-s", "format", "jpeg",
                 str(pdf_path), "--out", str(output_dir / "pdf_page.jpg")],
                capture_output=True, text=True, timeout=60,
                stdin=subprocess.DEVNULL
            )
            # sips 只能轉第一頁，較有限
            out = output_dir / "pdf_page.jpg"
            if out.exists():
                images.append(out)
        except Exception:
            pass
    except Exception as e:
        log(f"  ❌ PDF 處理失敗: {e}")
    return images


def analyze_slides(folder: Path) -> Path | None:
    """分析資料夾中的投影片內容，產出 slides_analysis.md"""
    slide_files = find_slide_files(folder)
    if not slide_files:
        return None

    output_file = folder / "slides_analysis.md"
    if output_file.exists():
        log(f"  ℹ️ slides_analysis.md 已存在，跳過分析")
        return output_file

    log(f"  🖼️ 開始分析 {len(slide_files)} 個投影片檔案")

    # 收集所有圖片（直接圖片 + 從 PPTX/PDF 提取的）
    all_images: list[Path] = []
    temp_dir = folder / "_temp_extracted"

    for sf in slide_files:
        ext = _ext(sf)
        if ext in IMAGE_EXTENSIONS:
            all_images.append(sf)
        elif ext == "pptx":
            temp_dir.mkdir(exist_ok=True)
            extracted = extract_images_from_pptx(sf, temp_dir)
            all_images.extend(extracted)
        elif ext == "pdf":
            temp_dir.mkdir(exist_ok=True)
            extracted = extract_images_from_pdf(sf, temp_dir)
            all_images.extend(extracted)

    if not all_images:
        log(f"  ⚠️ 無可分析的圖片")
        return None

    log(f"  🔍 共 {len(all_images)} 張圖片待分析")

    # 用 markitdown_helper 分析
    env = get_env()
    api_key = env.get("OPENAI_API_KEY", "")

    try:
        # 動態 import markitdown_helper（在 SKILL_DIR 同目錄下）
        if str(SKILL_DIR) not in sys.path:
            sys.path.insert(0, str(SKILL_DIR))
        from markitdown_helper import convert_images_to_markdown

        success, out_path, info = convert_images_to_markdown(
            image_paths=[str(p) for p in all_images],
            output_file=str(output_file),
            title=f"{folder.name} 投影片分析",
            use_llm=True,
            api_key=api_key,
            model=OPENAI_VISION
        )

        if success:
            log(f"  ✅ 投影片分析完成: {output_file.name} ({info.get('processed_images', 0)} 張)")
            return output_file
        else:
            log(f"  ❌ 投影片分析失敗: {info.get('error', 'unknown')}")
            return None

    except ImportError:
        log(f"  ⚠️ markitdown_helper 未找到，嘗試直接呼叫 OpenAI Vision...")
        # Fallback: 直接用 OpenAI Vision API
        return _analyze_slides_direct(all_images, output_file, api_key)
    finally:
        # 清理臨時提取的圖片
        if temp_dir.exists():
            shutil.rmtree(temp_dir, ignore_errors=True)


def _analyze_slides_direct(images: list[Path], output_file: Path, api_key: str) -> Path | None:
    """直接用 OpenAI Vision API 分析圖片（fallback）"""
    try:
        import base64
        from openai import OpenAI

        client = OpenAI(api_key=api_key)

        with open(output_file, "w", encoding="utf-8") as f:
            f.write(f"# {output_file.parent.name} 投影片分析\n\n")

            for i, img_path in enumerate(images):
                slide_num = i + 1
                f.write(f"## 投影片 {slide_num}\n\n")

                log(f"    分析投影片 {slide_num}/{len(images)}: {img_path.name}")

                with open(img_path, "rb") as img_file:
                    encoded = base64.b64encode(img_file.read()).decode("utf-8")

                response = client.chat.completions.create(
                    model=OPENAI_VISION,
                    messages=[
                        {"role": "system", "content": (
                            "你是一個醫學投影片分析專家。請識別並提取圖片中所有可見的"
                            "文本內容，同時分析圖表、表格和視覺元素。"
                            "以結構化的 Markdown 格式返回，使用繁體中文。"
                        )},
                        {"role": "user", "content": [
                            {"type": "text", "text": "請分析這張投影片並提取內容："},
                            {"type": "image_url", "image_url": {
                                "url": f"data:image/jpeg;base64,{encoded}"
                            }}
                        ]}
                    ],
                    max_completion_tokens=1500
                )

                f.write(f"{response.choices[0].message.content}\n\n---\n\n")

        log(f"  ✅ 投影片分析完成: {output_file.name}")
        return output_file

    except Exception as e:
        log(f"  ❌ 投影片分析失敗: {e}")
        return None


# ===== C. 合併處理流程 =====

def _find_existing_transcript(folder: Path) -> Path | None:
    """尋找已存在的逐字稿：先看資料夾內，再看 ~/audio/ 同名 .txt"""
    # 1. 資料夾內的 .txt
    txt_files = list(folder.glob("*.txt"))
    if txt_files:
        return txt_files[0]

    # 2. ~/audio/ 中與資料夾同名的 .txt（之前單獨處理留下的）
    parent_txt = folder.parent / f"{folder.name}.txt"
    if parent_txt.exists():
        # 複製到資料夾內，方便後續使用
        import shutil
        dest = folder / parent_txt.name
        shutil.copy2(parent_txt, dest)
        log(f"  📋 找到先前的逐字稿: {parent_txt.name} → 複製到資料夾")
        return dest

    return None


def process_lecture_folder(folder: Path, state: dict) -> bool:
    """處理一個演講資料夾：音訊 + 投影片 → 合併筆記

    支援兩種流程：
    1. 音訊+投影片同時放入 → 一次處理完
    2. 先單獨處理音訊（~/audio/name.m4a），之後建立 ~/audio/name/ 放投影片
       → 自動找到之前的逐字稿，分析投影片後合併
    """
    folder_name = folder.name
    folder_state = state["processed_folders"].get(folder_name, {})

    audio_files = find_audio_files(folder)
    slide_files = find_slide_files(folder)

    if not audio_files and not slide_files:
        return False

    # 檢查是否還有事要做
    needs_audio = bool(audio_files) and not folder_state.get("audio_done")
    needs_slides = bool(slide_files) and not folder_state.get("slides_done")
    # 可能需要合併：投影片已分析，但尚未合併，且有音訊或先前逐字稿可用
    needs_merge = (folder_state.get("slides_done") and not folder_state.get("merged")
                   and (bool(audio_files) or (folder.parent / f"{folder.name}.txt").exists()))
    if not needs_audio and not needs_slides and not needs_merge:
        return False

    log(f"📂 處理資料夾: {folder_name}")
    if audio_files:
        log(f"  🎵 音訊: {len(audio_files)} 個{'（已完成）' if folder_state.get('audio_done') else ''}")
    if slide_files:
        log(f"  🖼️ 投影片: {len(slide_files)} 個{'（已完成）' if folder_state.get('slides_done') else ''}")

    results = {"audio_done": False, "slides_done": False, "merged": False}

    # --- A. 處理音訊 ---
    transcript_path = None
    notes_path = None

    if audio_files and not folder_state.get("audio_done"):
        audio_file = audio_files[0]  # 取第一個音訊檔
        transcript_path = transcribe_audio(audio_file, folder)

        if transcript_path:
            results["audio_done"] = True

            # 先檢查有沒有 slides 分析可以合併
            slides_md = folder / "slides_analysis.md"
            if not slides_md.exists() and slide_files:
                slides_md = analyze_slides(folder)

            # 產生筆記（有 slides 就 merge，沒有就 transcript-only）
            notes_path = summarize_transcript(transcript_path, slides_md)
            if notes_path:
                convert_md_to_docx(notes_path)
                results["merged"] = slides_md is not None and slides_md.exists()

    elif folder_state.get("audio_done"):
        # 音訊已在此資料夾處理過
        transcript_path = _find_existing_transcript(folder)

    elif not audio_files:
        # 資料夾沒有音訊，但可能之前以單一檔案處理過
        # 例如：~/audio/pit-1.m4a 已處理 → 現在有 ~/audio/pit-1/ 放投影片
        transcript_path = _find_existing_transcript(folder)
        if transcript_path:
            log(f"  📋 使用先前的逐字稿: {transcript_path.name}")
            results["audio_done"] = True  # 標記為已完成（用既有的）

    # --- B. 處理投影片（如果有投影片尚未處理）---
    if slide_files and not folder_state.get("slides_done"):
        slides_md = folder / "slides_analysis.md"
        if not slides_md.exists():
            slides_md = analyze_slides(folder)
            results["slides_done"] = slides_md is not None

            # 有逐字稿 + 投影片分析 → 合併筆記
            if slides_md and transcript_path and not folder_state.get("merged"):
                log(f"  🔄 逐字稿+投影片到齊，產生合併筆記...")
                notes_path = summarize_transcript(transcript_path, slides_md)
                if notes_path:
                    convert_md_to_docx(notes_path)
                    results["merged"] = True
        else:
            results["slides_done"] = True

    # 更新狀態
    folder_state.update({
        "audio_done": folder_state.get("audio_done") or results["audio_done"],
        "slides_done": folder_state.get("slides_done") or results["slides_done"],
        "merged": folder_state.get("merged") or results["merged"],
        "last_processed": datetime.now().isoformat(),
    })
    state["processed_folders"][folder_name] = folder_state
    save_state(state)

    # Telegram 通知
    if results["audio_done"] or results["slides_done"]:
        status = []
        if results["audio_done"]:
            status.append("🎙️ 轉錄完成")
        if results["slides_done"]:
            status.append("🖼️ 投影片分析完成")
        if results["merged"]:
            status.append("📝 合併筆記完成")
        send_telegram(
            f"📂 *{folder_name}*\n" + "\n".join(status)
        )

    return results["audio_done"] or results["slides_done"]


def process_audio_only(audio_path: Path, state: dict) -> bool:
    """處理 ~/audio/ 中的單一音訊檔（僅轉錄 + 筆記）

    轉錄產出的 .txt 保留在 ~/audio/，方便之後建子資料夾時合併使用。
    音訊檔本身移至 ~/audio/done/。
    """
    if str(audio_path) in state.get("processed_audio", []):
        return False

    log(f"🎵 處理音訊: {audio_path.name}")

    output_dir = audio_path.parent
    transcript_path = transcribe_audio(audio_path, output_dir)

    if transcript_path:
        notes_path = summarize_transcript(transcript_path)
        if notes_path:
            convert_md_to_docx(notes_path)

        # 音訊檔移到 done（但 .txt, .md, .docx 留在 ~/audio/ 供後續合併）
        try:
            done_path = DONE_DIR / audio_path.name
            audio_path.rename(done_path)
            log(f"  📁 音訊已移至 done/（逐字稿保留供後續合併）")
        except Exception:
            pass

        state.setdefault("processed_audio", []).append(str(audio_path))
        save_state(state)

        send_telegram(
            f"🎵 *音訊處理完成*\n"
            f"📎 {audio_path.name}\n"
            f"📝 筆記已產生\n"
            f"💡 如需合併投影片，建立 ~/audio/{audio_path.stem}/ 放入投影片即可"
        )
        return True

    return False


def scan_folders():
    """掃描 ~/audio/ 資料夾：檔案=單純轉錄，子資料夾=音訊+投影片合併"""
    state = load_state()
    has_work = False

    if not AUDIO_DIR.exists():
        return

    # 1. 掃描子資料夾（音訊+投影片 → 合併筆記）
    subfolders = sorted([
        d for d in AUDIO_DIR.iterdir()
        if d.is_dir() and not d.name.startswith(".") and d.name != "done"
    ])
    for folder in subfolders:
        try:
            if process_lecture_folder(folder, state):
                has_work = True
        except Exception as e:
            log(f"  ❌ 處理失敗 {folder.name}: {e}")
            import traceback
            traceback.print_exc()

    # 2. 掃描單一音訊檔（僅轉錄+筆記）
    audio_files = sorted([
        f for f in AUDIO_DIR.iterdir()
        if f.is_file() and _ext(f) in AUDIO_EXTENSIONS
        and not f.name.startswith(".")
        and str(f) not in state.get("processed_audio", [])
    ])
    for af in audio_files:
        try:
            if process_audio_only(af, state):
                has_work = True
        except Exception as e:
            log(f"  ❌ 處理失敗 {af.name}: {e}")

    if not has_work and not subfolders and not audio_files:
        log(f"📂 掃描 {AUDIO_DIR}... 無新檔案")


def main():
    parser = argparse.ArgumentParser(description="資料夾監控處理器")
    parser.add_argument("--daemon", action="store_true", help="持續監控模式")
    parser.add_argument("--interval", type=int, default=POLL_INTERVAL, help="輪詢間隔（秒）")
    args = parser.parse_args()

    setup_dirs()

    if args.daemon:
        log("🚀 Folder Processor 啟動（daemon 模式）")
        log(f"  📂 監控: {AUDIO_DIR}")
        log(f"  📁 完成: {DONE_DIR}")
        log(f"  ⏱️ 間隔: {args.interval}s")
        log(f"  📋 規則: 檔案→轉錄, 子資料夾→轉錄+投影片+合併")
        while True:
            try:
                scan_folders()
            except Exception as e:
                log(f"❌ 掃描錯誤: {e}")
            time.sleep(args.interval)
    else:
        scan_folders()


if __name__ == "__main__":
    main()
